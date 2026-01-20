"""
PC047组会PPT生成脚本 v5
核心叙述：CagA感染如何重塑肠道微生物组
新增：02b GO/PFAM验证 + 03c共现网络分析
优化：7分钟演讲适配
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'green': RGBColor(0x4C, 0xAF, 0x50),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'orange': RGBColor(0xFF, 0x98, 0x00),
    'red': RGBColor(0xE5, 0x39, 0x35),
    'purple': RGBColor(0x7B, 0x1F, 0xA2),
    'teal': RGBColor(0x00, 0x96, 0x88),
}


def add_page_number(slide, page_num, slide_width, slide_height):
    """为幻灯片添加页码（右下角，蓝色加粗）"""
    page_box = slide.shapes.add_textbox(
        slide_width - Inches(1.0),
        slide_height - Inches(0.6),
        Inches(0.8),
        Inches(0.4)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.RIGHT


def add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False):
    """为整个演示文稿添加页码"""
    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, 1):
        if skip_first and idx == 1:
            continue
        if skip_last and idx == total_slides:
            continue
        add_page_number(slide, idx, prs.slide_width, prs.slide_height)


def add_header(slide, prs, title):
    """添加统一标题栏"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_blue']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']


def add_box_with_text(slide, left, top, width, height, text, fill_color, text_color=None, font_size=14, bold=False):
    """添加带文字的方框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS['dark_gray']
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color or COLORS['dark_gray']
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_title_slide(prs, title, subtitle, date, presenter):
    """添加封面幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_blue']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{date}\n{presenter}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_background_slide_1(prs, lang='cn'):
    """背景页1：CagA与肠道肿瘤"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：CagA与肠道肿瘤" if lang == 'cn' else "Background: CagA and Intestinal Tumors"
    add_header(slide, prs, title)

    # 左侧：H. pylori感染流程
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "H. pylori CagA 致病机制" if lang == 'cn' else "H. pylori CagA Pathogenesis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    boxes_left = [
        (0.5, 1.7, 1.8, 0.5, "H. pylori\n感染" if lang == 'cn' else "H. pylori\nInfection", COLORS['light_blue']),
        (2.6, 1.7, 1.8, 0.5, "CagA蛋白\n注入宿主细胞" if lang == 'cn' else "CagA Injection\ninto Host Cell", COLORS['light_blue']),
        (0.5, 2.5, 1.8, 0.6, "干扰Wnt/NF-κB\n信号通路" if lang == 'cn' else "Disrupt Wnt/\nNF-κB Pathways", COLORS['orange']),
        (2.6, 2.5, 1.8, 0.6, "肠道菌群\n失调" if lang == 'cn' else "Gut Microbiota\nDysbiosis", COLORS['orange']),
        (1.55, 3.4, 2.0, 0.5, "促进肿瘤发生" if lang == 'cn' else "Tumor Promotion", COLORS['red']),
    ]

    for left, top, width, height, text, color in boxes_left:
        text_color = COLORS['white'] if color == COLORS['red'] else COLORS['dark_gray']
        add_box_with_text(slide, left, top, width, height, text, color, text_color, font_size=11)

    # 箭头
    line1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.35), Inches(1.85), Inches(0.2), Inches(0.15))
    line1.fill.solid()
    line1.fill.fore_color.rgb = COLORS['dark_gray']
    line1.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(2.25), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(3.15), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    # 右侧：关键文献支持
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "关键文献支持" if lang == 'cn' else "Key Literature Support"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    refs = [
        ("Jones et al. 2017, PLoS Pathog", "果蝇模型：CagA单独表达即可\n诱导菌群失调和上皮过度增殖" if lang == 'cn' else "Drosophila: CagA alone induces\ndysbiosis and epithelial proliferation"),
        ("Cui et al. 2025, BMC Gastro", "临床研究：CagA+菌株通过\n调节肠道菌群影响结直肠病变" if lang == 'cn' else "Clinical: CagA+ strains affect\ncolorectal lesions via gut microbiota"),
        ("Ding et al. 2025, Cell Host", "噬菌体可靶向清除促肿瘤\n细菌并恢复化疗敏感性" if lang == 'cn' else "Phage can eliminate tumor-\npromoting bacteria"),
    ]

    y_pos = 1.7
    for ref, desc in refs:
        ref_box = slide.shapes.add_textbox(Inches(5.0), Inches(y_pos), Inches(4.5), Inches(0.3))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + ref
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary_blue']

        desc_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.3), Inches(4.3), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['dark_gray']

        y_pos += 0.85

    # 底部：核心问题框
    question_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
    )
    question_box.fill.solid()
    question_box.fill.fore_color.rgb = COLORS['dark_blue']
    question_box.line.fill.background()

    q_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(8.6), Inches(0.6))
    tf = q_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "核心问题：CagA抗原清除后，肠道内什么因素维持CD8+ T细胞的持续激活？"
    else:
        p.text = "Central Question: What maintains CD8+ T-cell activation after CagA clearance?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """首先介绍研究背景。（30秒）

CagA蛋白通过IV型分泌系统注入宿主细胞，干扰Wnt和NF-κB信号通路，同时导致肠道菌群失调，最终促进肿瘤发生。

核心科学问题：CagA抗原清除后，肠道内是什么因素维持着CD8+ T细胞的持续激活？"""
    else:
        notes = """Research background. (30s)

CagA protein disrupts Wnt and NF-κB pathways, causes gut microbiota dysbiosis, promoting tumor development.

Central question: What maintains CD8+ T-cell activation after CagA clearance?"""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_background_slide_2(prs, lang='cn'):
    """背景页2：Apc突变与G×E交互"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：Apc突变与遗传易感性" if lang == 'cn' else "Background: Apc Mutation and Genetic Susceptibility"
    add_header(slide, prs, title)

    # 左侧：Apc-Wnt通路示意
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变 → Wnt通路持续激活" if lang == 'cn' else "Apc Mutation → Constitutive Wnt Activation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 正常状态
    normal_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(2), Inches(0.3))
    tf = normal_label.text_frame
    p = tf.paragraphs[0]
    p.text = "正常状态" if lang == 'cn' else "Normal"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    add_box_with_text(slide, 0.5, 1.9, 1.0, 0.4, "APC", COLORS['green'], COLORS['white'], 12, True)
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.0), Inches(0.3), Inches(0.2))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = COLORS['dark_gray']
    arr1.line.fill.background()
    add_box_with_text(slide, 1.9, 1.9, 1.2, 0.4, "β-catenin\n降解" if lang == 'cn' else "β-catenin\nDegradation", COLORS['light_gray'], COLORS['dark_gray'], 10)
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.0), Inches(0.3), Inches(0.2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = COLORS['dark_gray']
    arr2.line.fill.background()
    add_box_with_text(slide, 3.5, 1.9, 1.0, 0.4, "稳态" if lang == 'cn' else "Homeostasis", COLORS['green'], COLORS['white'], 11, True)

    # 突变状态
    mut_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(0.3))
    tf = mut_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变" if lang == 'cn' else "Apc Mutant"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['red']

    add_box_with_text(slide, 0.5, 2.8, 1.0, 0.4, "APC ✗", COLORS['red'], COLORS['white'], 12, True)
    arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.9), Inches(0.3), Inches(0.2))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = COLORS['dark_gray']
    arr3.line.fill.background()
    add_box_with_text(slide, 1.9, 2.8, 1.2, 0.4, "β-catenin\n累积" if lang == 'cn' else "β-catenin\nAccumulation", COLORS['orange'], COLORS['dark_gray'], 10)
    arr4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.9), Inches(0.3), Inches(0.2))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = COLORS['dark_gray']
    arr4.line.fill.background()
    add_box_with_text(slide, 3.5, 2.8, 1.0, 0.4, "肿瘤" if lang == 'cn' else "Tumor", COLORS['red'], COLORS['white'], 11, True)

    # 右侧：G×E交互概念图
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "G×E 交互作用假说" if lang == 'cn' else "G×E Interaction Hypothesis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 2x2矩阵标签
    col_labels = [("ApcWT", 6.0), ("ApcMUT", 7.8)]
    row_labels = [("CagA-", 2.0), ("CagA+", 2.8)]

    for label, x in col_labels:
        lb = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(1.5), Inches(0.3))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

    for label, y in row_labels:
        lb = slide.shapes.add_textbox(Inches(5.0), Inches(y), Inches(0.8), Inches(0.5))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']

    # 矩阵格子
    matrix_data = [
        (6.0, 1.95, "正常\n菌群" if lang == 'cn' else "Normal\nMicrobiota", COLORS['green']),
        (7.8, 1.95, "轻度\n失调" if lang == 'cn' else "Mild\nDysbiosis", COLORS['orange']),
        (6.0, 2.75, "无效应\nP=0.387" if lang == 'cn' else "No Effect\nP=0.387", COLORS['light_gray']),
        (7.8, 2.75, "显著重塑\nP=0.039" if lang == 'cn' else "Significant\nP=0.039", COLORS['red']),
    ]

    for x, y, text, color in matrix_data:
        text_color = COLORS['white'] if color in [COLORS['red'], COLORS['green']] else COLORS['dark_gray']
        add_box_with_text(slide, x, y, 1.5, 0.6, text, color, text_color, 11, True)

    # 结论框
    conclusion = slide.shapes.add_textbox(Inches(5.0), Inches(3.5), Inches(4.5), Inches(0.6))
    tf = conclusion.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "→ CagA效应需要Apc突变的易感背景"
    else:
        p.text = "→ CagA effect requires Apc-mutant susceptible background"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 底部：实验设计简图
    design_title = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(4.5), Inches(0.4))
    tf = design_title.text_frame
    p = tf.paragraphs[0]
    p.text = "实验设计：2×3因子" if lang == 'cn' else "Experimental Design: 2×3 Factorial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    groups = [
        (0.5, 3.9, "ApcWT", COLORS['light_blue']),
        (2.2, 3.9, "ApcMUT", COLORS['orange']),
    ]
    for x, y, text, color in groups:
        add_box_with_text(slide, x, y, 1.5, 0.4, text, color, COLORS['dark_gray'], 12, True)

    infections = ["Ctrl", "HpKO", "HpWT"]
    for i, inf in enumerate(infections):
        x = 0.5 + i * 1.4
        add_box_with_text(slide, x, 4.4, 1.2, 0.35, inf, COLORS['light_gray'], COLORS['dark_gray'], 10)

    core_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(4.85), Inches(1.7), Inches(0.35)
    )
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = COLORS['red']
    core_box.line.fill.background()

    core_text = slide.shapes.add_textbox(Inches(2.8), Inches(4.87), Inches(1.7), Inches(0.3))
    tf = core_text.text_frame
    p = tf.paragraphs[0]
    p.text = "核心比较" if lang == 'cn' else "Core Comparison"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """Apc突变与G×E交互作用。（30秒）

CagA感染的效应高度依赖宿主遗传背景。野生型背景下CagA几乎无效应（P=0.387），但在Apc突变背景下显著重塑菌群（P=0.039）。

核心比较是Apc突变背景下CagA阳性和阴性组。"""
    else:
        notes = """Apc mutation and G×E interaction. (30s)

CagA effect is highly dependent on host genetics. No effect in wild-type (P=0.387), but significant restructuring in Apc-mutant (P=0.039).

Core comparison: CagA+ vs CagA- under Apc-mutant background."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_hypothesis_slide(prs, lang='cn'):
    """假说图解页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "核心假说：Functional Footprint" if lang == 'cn' else "Central Hypothesis: Functional Footprint"
    add_header(slide, prs, title)

    steps = [
        ("1", "CagA感染\n重塑菌群" if lang == 'cn' else "CagA Infection\nRestructures Microbiota", COLORS['primary_blue'], 0.5),
        ("2", "菌群产生\n特异性分子" if lang == 'cn' else "Microbiota Produces\nSpecific Molecules", COLORS['teal'], 2.7),
        ("3", "持续激活\nCD8+ T细胞" if lang == 'cn' else "Persistent CD8+\nT-cell Activation", COLORS['orange'], 4.9),
        ("4", "促进\n肿瘤发生" if lang == 'cn' else "Tumor\nPromotion", COLORS['red'], 7.1),
    ]

    for num, text, color, x in steps:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(1.4), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(x), Inches(1.45), Inches(0.4), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        add_box_with_text(slide, x, 1.9, 1.8, 0.7, text, COLORS['light_gray'], color, 12, True)

        if x < 7:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.85), Inches(2.15), Inches(0.7), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS['dark_gray']
            arr.line.fill.background()

    # 验证状态
    status_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.4))
    tf = status_title.text_frame
    p = tf.paragraphs[0]
    p.text = "本研究验证范围" if lang == 'cn' else "Scope of This Study"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    verified_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(4.2), Inches(1.5)
    )
    verified_box.fill.solid()
    verified_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    verified_box.line.color.rgb = COLORS['green']
    verified_box.line.width = Pt(2)

    v_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.8), Inches(0.3))
    tf = v_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ 已验证" if lang == 'cn' else "✓ Verified"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    v_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(3.8), Inches(1))
    tf = v_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 菌群结构显著重塑 (P=0.012)\n• G×E交互作用确认 (9/9验证)\n• 功能冗余 (GO/PFAM全验证)\n• 网络结构瓦解 (模块度-62%)"
    else:
        p.text = "• Microbiota restructured (P=0.012)\n• G×E interaction confirmed (9/9 verified)\n• Functional redundancy (GO/PFAM verified)\n• Network collapse (Modularity -62%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    pending_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.3), Inches(4.5), Inches(1.5)
    )
    pending_box.fill.solid()
    pending_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    pending_box.line.color.rgb = COLORS['orange']
    pending_box.line.width = Pt(2)

    p_title = slide.shapes.add_textbox(Inches(5.2), Inches(3.4), Inches(4.1), Inches(0.3))
    tf = p_title.text_frame
    p = tf.paragraphs[0]
    p.text = "○ 待验证" if lang == 'cn' else "○ Pending"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['orange']

    p_content = slide.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4.1), Inches(1))
    tf = p_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 代谢组学：特定代谢物鉴定\n• 转录组学：功能基因表达验证\n• 整合肿瘤/T细胞表型数据"
    else:
        p.text = "• Metabolomics: Specific metabolite ID\n• Transcriptomics: Gene expression validation\n• Integration with tumor/T-cell phenotypes"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    if lang == 'cn':
        notes = """这是核心假说——Functional Footprint假说。（30秒）

假说逻辑链：CagA感染重塑菌群 → 产生特异性分子 → 持续激活CD8+ T细胞 → 促进肿瘤发生。

本次验证的前半部分：菌群重塑确认、G×E交互验证、功能冗余（GO/PFAM全层级验证）、网络结构瓦解（模块度降低62%）。"""
    else:
        notes = """Central hypothesis - Functional Footprint. (30s)

Logical chain: CagA restructures microbiota → produces specific molecules → persistent T-cell activation → tumor promotion.

Verified: microbiota restructuring, G×E interaction, functional redundancy (GO/PFAM verified), network collapse (modularity -62%)."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_analysis_pipeline_slide(prs, lang='cn'):
    """分析流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "分析流程概览" if lang == 'cn' else "Analysis Pipeline Overview"
    add_header(slide, prs, title)

    add_box_with_text(slide, 3.8, 1.3, 2.4, 0.5,
                      "盲肠样本 (n=27)" if lang == 'cn' else "Cecum Samples (n=27)",
                      COLORS['primary_blue'], COLORS['white'], 14, True)

    modules = [
        {
            'title': "01 物种组成分析" if lang == 'cn' else "01 Species Analysis",
            'tool': "Kraken2/Bracken",
            'focus': "群落结构\nG×E交互" if lang == 'cn' else "Community\nG×E Interaction",
            'result': "5/9显著" if lang == 'cn' else "5/9 Sig.",
            'x': 0.3,
            'color': COLORS['primary_blue'],
        },
        {
            'title': "02 功能潜力分析" if lang == 'cn' else "02 Functional Analysis",
            'tool': "HUMAnN4",
            'focus': "Pathway→KO\n→GO→PFAM" if lang == 'cn' else "Pathway→KO\n→GO→PFAM",
            'result': "功能冗余" if lang == 'cn' else "Redundancy",
            'x': 3.5,
            'color': COLORS['teal'],
        },
        {
            'title': "03 网络/噬菌体" if lang == 'cn' else "03 Network/Phage",
            'tool': "igraph/Lyrebird",
            'focus': "共现网络\n细菌-噬菌体" if lang == 'cn' else "Co-occurrence\nBacteria-Phage",
            'result': "模块度-62%" if lang == 'cn' else "Mod -62%",
            'x': 6.7,
            'color': COLORS['purple'],
        },
    ]

    for mod in modules:
        x = mod['x']

        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.0), Inches(0.45)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = mod['color']
        title_box.line.fill.background()

        t_text = slide.shapes.add_textbox(Inches(x), Inches(2.05), Inches(3.0), Inches(0.4))
        tf = t_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['title']
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        tool_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.0), Inches(0.35))
        tf = tool_box.text_frame
        p = tf.paragraphs[0]
        p.text = mod['tool']
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        focus_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.3), Inches(2.9), Inches(2.4), Inches(0.7)
        )
        focus_box.fill.solid()
        focus_box.fill.fore_color.rgb = COLORS['light_gray']
        focus_box.line.color.rgb = mod['color']

        f_text = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.95), Inches(2.4), Inches(0.65))
        tf = f_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mod['focus']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        result_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.6), Inches(3.7), Inches(1.8), Inches(0.4)
        )
        result_box.fill.solid()
        result_box.fill.fore_color.rgb = mod['color']
        result_box.line.fill.background()

        r_text = slide.shapes.add_textbox(Inches(x + 0.6), Inches(3.73), Inches(1.8), Inches(0.35))
        tf = r_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['result']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    for start_x in [1.8, 3.3, 5.0, 6.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(start_x + 2.9), Inches(1.85), Inches(0.2), Inches(0.12))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(9), Inches(0.9)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(8.6), Inches(0.75))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "整合结论：CagA重塑细菌菌群 → 功能冗余(全层级验证) → 网络结构瓦解 → 噬菌体协同"
    else:
        p.text = "Integration: CagA restructures bacteria → Functional redundancy (all levels) → Network collapse → Phage coordination"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """分析流程概览。（30秒）

三大分析模块：物种组成分析（9个比较5个显著）、功能潜力分析（从Pathway到PFAM全层级功能冗余）、网络/噬菌体分析（模块度降低62%）。

整合结论：CagA重塑细菌菌群，功能冗余但网络结构瓦解，噬菌体协同变化。"""
    else:
        notes = """Analysis pipeline overview. (30s)

Three modules: species composition (5/9 significant), functional potential (all levels redundant), network/phage (modularity -62%).

Integration: CagA restructures bacteria, functional redundancy with network collapse, phage coordination."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_content_slide(prs, title, bullets, notes="", image_path=None):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    if image_path and os.path.exists(image_path):
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(4))
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if '**' in bullet:
            parts = bullet.split('**')
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(18)
                if j % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = COLORS['green']
                else:
                    run.font.color.rgb = COLORS['dark_gray']
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark_gray']

        p.space_before = Pt(8)
        p.level = 0

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_bacteria_vs_virus_slide(prs, lang='cn'):
    """新增：细菌vs病毒对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果2：CagA的作用靶点——细菌而非病毒" if lang == 'cn' else "Result 2: CagA Targets Bacteria, Not Viruses"
    add_header(slide, prs, title)

    # 问题引入
    question = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
    tf = question.text_frame
    p = tf.paragraphs[0]
    p.text = "问：CagA重塑的是什么？细菌？还是病毒？" if lang == 'cn' else "Q: What does CagA reshape? Bacteria or Viruses?"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['dark_blue']

    # 左侧：Standard细菌
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Standard (细菌)" if lang == 'cn' else "Standard (Bacteria)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 细菌结果框
    bacteria_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.3), Inches(4), Inches(2.0)
    )
    bacteria_box.fill.solid()
    bacteria_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    bacteria_box.line.color.rgb = COLORS['primary_blue']
    bacteria_box.line.width = Pt(2)

    bacteria_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = bacteria_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：5/9 (55.6%)"
    else:
        p.text = "• FDR Significant: 5/9 (55.6%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    else:
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.387"
    else:
        p.text = "• CagA@ApcWT: P=0.387"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ G×E交互完全验证"
    else:
        p.text = "→ G×E Fully Verified"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 右侧：Virus病毒
    right_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Virus (病毒)" if lang == 'cn' else "Virus (Viruses)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 病毒结果框
    virus_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.3), Inches(4), Inches(2.0)
    )
    virus_box.fill.solid()
    virus_box.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
    virus_box.line.color.rgb = COLORS['purple']
    virus_box.line.width = Pt(2)

    virus_content = slide.shapes.add_textbox(Inches(5.5), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = virus_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：0/9 (0%)"
    else:
        p.text = "• FDR Significant: 0/9 (0%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.604"
    else:
        p.text = "• CagA@ApcMUT: P=0.604"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.978"
    else:
        p.text = "• CagA@ApcWT: P=0.978"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ 无显著效应"
    else:
        p.text = "→ No Significant Effect"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 中间对比符号
    vs_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.55), Inches(3.0), Inches(0.7), Inches(0.7)
    )
    vs_box.fill.solid()
    vs_box.fill.fore_color.rgb = COLORS['dark_gray']
    vs_box.line.fill.background()

    vs_text = slide.shapes.add_textbox(Inches(4.55), Inches(3.15), Inches(0.7), Inches(0.4))
    tf = vs_text.text_frame
    p = tf.paragraphs[0]
    p.text = "VS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(0.5))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：CagA主要通过重塑细菌群落影响肠道微环境，而非直接作用于病毒组"
    else:
        p.text = "Conclusion: CagA reshapes gut environment mainly via bacteria, not directly via viruses"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """CagA到底影响什么？（30秒）

细菌9个比较5个显著，核心CagA效应P=0.039，完全验证G×E交互。

病毒9个比较0个显著，核心CagA比较P=0.604。

结论：CagA主要通过重塑细菌群落影响肠道微环境。"""
    else:
        notes = """What does CagA affect? (30s)

Bacteria: 5/9 significant, core CagA effect P=0.039, G×E verified.

Viruses: 0/9 significant, core comparison P=0.604.

Conclusion: CagA reshapes gut environment mainly via bacteria."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_functional_redundancy_slide(prs, lang='cn'):
    """新增：功能冗余全层级验证页（含GO/PFAM）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果3：功能冗余——全层级验证" if lang == 'cn' else "Result 3: Functional Redundancy — All Levels Verified"
    add_header(slide, prs, title)

    # 左侧：功能层级验证表
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HUMAnN4 功能层级验证" if lang == 'cn' else "HUMAnN4 Functional Hierarchy"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 表格数据
    table_data = [
        ("层级", "特征数", "P<0.05", "FDR<0.05", "结论"),
        ("Pathway", "~500", "多个", "0", "✓ 冗余"),
        ("KO", "~6,000", "多个", "0", "✓ 冗余"),
        ("GO", "~13,000", "1,577", "0", "✓ 冗余"),
        ("PFAM", "~7,600", "849", "0", "✓ 冗余"),
    ]

    # 创建表格
    table_top = 1.7
    row_height = 0.4
    col_widths = [1.2, 1.1, 0.9, 1.0, 0.9]

    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            x = 0.4 + sum(col_widths[:col_idx])
            y = table_top + row_idx * row_height

            if row_idx == 0:
                # 表头
                cell_bg = COLORS['primary_blue']
                text_color = COLORS['white']
                font_bold = True
            elif "✓" in str(cell_text):
                cell_bg = RGBColor(0xE8, 0xF5, 0xE9)
                text_color = COLORS['green']
                font_bold = True
            else:
                cell_bg = COLORS['light_gray'] if row_idx % 2 == 0 else COLORS['white']
                text_color = COLORS['dark_gray']
                font_bold = False

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[col_idx]), Inches(row_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = cell_bg
            cell.line.color.rgb = COLORS['dark_gray']
            cell.line.width = Pt(0.5)

            cell_text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(col_widths[col_idx]), Inches(0.3))
            tf = cell_text_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(10)
            p.font.bold = font_bold
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

    # 右侧：Driver更替示意
    right_title = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "静默更替：Driver物种改变" if lang == 'cn' else "Silent Shift: Driver Species Change"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # Driver更替框
    driver_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.7), Inches(4), Inches(2.3)
    )
    driver_box.fill.solid()
    driver_box.fill.fore_color.rgb = RGBColor(0xFB, 0xE9, 0xE7)
    driver_box.line.color.rgb = COLORS['orange']
    driver_box.line.width = Pt(2)

    driver_content = slide.shapes.add_textbox(Inches(5.7), Inches(1.8), Inches(3.6), Inches(2.1))
    tf = driver_content.text_frame
    tf.word_wrap = True

    if lang == 'cn':
        lines = [
            "功能总量不变，但执行者改变：",
            "",
            "益生菌 ↓",
            "  • L. johnsonii: -50%",
            "  • F. prausnitzii: ↓",
            "",
            "炎症相关菌 ↑",
            "  • M. schaedleri: +57%",
            "  • Enterocloster: ↑"
        ]
    else:
        lines = [
            "Total function unchanged, executors changed:",
            "",
            "Probiotics ↓",
            "  • L. johnsonii: -50%",
            "  • F. prausnitzii: ↓",
            "",
            "Inflammation-associated ↑",
            "  • M. schaedleri: +57%",
            "  • Enterocloster: ↑"
        ]

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "↓" in line and "益生菌" in line:
            p.font.color.rgb = COLORS['green']
            p.font.bold = True
        elif "↑" in line and ("炎症" in line or "Inflammation" in line):
            p.font.color.rgb = COLORS['red']
            p.font.bold = True
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(1.0)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(8.6), Inches(0.8))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：从Pathway→KO→GO→PFAM全层级FDR校正后均无显著差异\n功能冗余得到全面验证，但功能执行者已静默更替"
    else:
        p.text = "Conclusion: No FDR-significant differences at any level (Pathway→KO→GO→PFAM)\nFunctional redundancy fully verified, but executors silently shifted"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """功能冗余全层级验证。（45秒）

从粗到细的四个功能层面——Pathway、KO、GO、PFAM——FDR校正后均无显著差异，功能冗余得到全面验证。

但分层分析揭示"静默更替"：功能总量不变，执行者已换。益生菌L. johnsonii下降50%，炎症相关菌M. schaedleri上升57%。

这解释了为什么物种变化但功能稳定——不同物种可以执行相同功能。"""
    else:
        notes = """Functional redundancy at all levels. (45s)

From Pathway to KO to GO to PFAM — no FDR-significant differences at any level. Functional redundancy fully verified.

But stratified analysis reveals "Silent Shift": total function unchanged, executors changed. L. johnsonii -50%, M. schaedleri +57%.

This explains why species change but function remains stable — different species can perform the same functions."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_network_analysis_slide(prs, lang='cn'):
    """新增：共现网络分析页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果4：共现网络——模块结构瓦解" if lang == 'cn' else "Result 4: Co-occurrence Network — Module Collapse"
    add_header(slide, prs, title)

    # 左侧：网络拓扑对比
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "网络拓扑变化 (Top 100物种)" if lang == 'cn' else "Network Topology Change (Top 100 Species)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # HpKO (对照) 框
    hpko_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.75), Inches(2.2), Inches(1.6)
    )
    hpko_box.fill.solid()
    hpko_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    hpko_box.line.color.rgb = COLORS['green']
    hpko_box.line.width = Pt(2)

    hpko_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(2), Inches(0.3))
    tf = hpko_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HpKO (对照)" if lang == 'cn' else "HpKO (Control)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    hpko_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(2), Inches(1.2))
    tf = hpko_content.text_frame
    tf.word_wrap = True
    if lang == 'cn':
        lines = ["节点: 80", "边数: 1,280", "模块度: 0.468", "→ 清晰模块结构"]
    else:
        lines = ["Nodes: 80", "Edges: 1,280", "Modularity: 0.468", "→ Clear modules"]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "模块度" in line or "Modularity" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['green']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.7), Inches(2.4), Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS['red']
    arrow.line.fill.background()

    # HpWT (CagA+) 框
    hpwt_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.75), Inches(2.2), Inches(1.6)
    )
    hpwt_box.fill.solid()
    hpwt_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    hpwt_box.line.color.rgb = COLORS['red']
    hpwt_box.line.width = Pt(2)

    hpwt_title = slide.shapes.add_textbox(Inches(3.3), Inches(1.8), Inches(2), Inches(0.3))
    tf = hpwt_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HpWT (CagA+)" if lang == 'cn' else "HpWT (CagA+)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['red']

    hpwt_content = slide.shapes.add_textbox(Inches(3.3), Inches(2.1), Inches(2), Inches(1.2))
    tf = hpwt_content.text_frame
    tf.word_wrap = True
    if lang == 'cn':
        lines = ["节点: 100", "边数: 2,190 (+71%)", "模块度: 0.177", "→ 模块结构瓦解"]
    else:
        lines = ["Nodes: 100", "Edges: 2,190 (+71%)", "Modularity: 0.177", "→ Module collapse"]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "模块度" in line or "Modularity" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['red']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 模块度变化强调框
    mod_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.5), Inches(5.0), Inches(0.7)
    )
    mod_box.fill.solid()
    mod_box.fill.fore_color.rgb = COLORS['orange']
    mod_box.line.fill.background()

    mod_text = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(4.8), Inches(0.5))
    tf = mod_text.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "模块度: 0.468 → 0.177 (降低62%)"
    else:
        p.text = "Modularity: 0.468 → 0.177 (-62%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 右侧：Hub物种
    right_title = slide.shapes.add_textbox(Inches(5.8), Inches(1.3), Inches(3.7), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Hub物种 (网络枢纽)" if lang == 'cn' else "Hub Species (Network Hubs)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    hub_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.7), Inches(3.7), Inches(2.5)
    )
    hub_box.fill.solid()
    hub_box.fill.fore_color.rgb = COLORS['light_gray']
    hub_box.line.color.rgb = COLORS['purple']
    hub_box.line.width = Pt(1)

    hub_content = slide.shapes.add_textbox(Inches(5.95), Inches(1.8), Inches(3.4), Inches(2.3))
    tf = hub_content.text_frame
    tf.word_wrap = True

    if lang == 'cn':
        hub_lines = [
            "Top 5 Hub物种（按度排序）：",
            "• Blautia obeum (度=44)",
            "• Roseburia hominis (度=42)",
            "• Hungatella hathewayi (度=42)",
            "",
            "特征：主要为厌氧梭菌类",
            "功能：SCFA生产者",
            "",
            "桥梁物种：",
            "• P. distasonis (介数=402)"
        ]
    else:
        hub_lines = [
            "Top 5 Hub Species (by degree):",
            "• Blautia obeum (d=44)",
            "• Roseburia hominis (d=42)",
            "• Hungatella hathewayi (d=42)",
            "",
            "Feature: Mainly anaerobic Clostridia",
            "Function: SCFA producers",
            "",
            "Bridge species:",
            "• P. distasonis (btw=402)"
        ]

    for i, line in enumerate(hub_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        if "Top" in line or "桥梁" in line or "Bridge" in line or "特征" in line or "Feature" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['purple']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.35), Inches(9), Inches(0.85)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.45), Inches(8.6), Inches(0.7))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：CagA导致网络模块度降低62%，物种间界限模糊\n→ 生态位重叠增加 → 解释功能冗余机制"
    else:
        p.text = "Conclusion: CagA causes 62% modularity drop, blurred species boundaries\n→ Increased niche overlap → Explains functional redundancy mechanism"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """共现网络分析——模块结构瓦解。（45秒）

对照组HpKO：80个节点，1280条边，模块度0.468，有清晰的功能模块。

CagA感染组HpWT：100个节点，2190条边增加71%，模块度降到0.177。

模块度降低62%意味着物种间界限模糊，生态位重叠增加。这从网络层面解释了功能冗余的机制——原本各司其职的物种现在功能重叠了。

Hub物种主要是SCFA生产者，桥梁物种P. distasonis可能是关键免疫调节点。"""
    else:
        notes = """Co-occurrence network analysis — module collapse. (45s)

Control HpKO: 80 nodes, 1,280 edges, modularity 0.468, clear functional modules.

CagA+ HpWT: 100 nodes, 2,190 edges (+71%), modularity dropped to 0.177.

62% modularity drop means blurred species boundaries, increased niche overlap. This explains functional redundancy at network level — species that used to specialize now have overlapping functions.

Hub species are mainly SCFA producers; bridge species P. distasonis may be key immune modulation point."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_phage_coordination_slide(prs, lang='cn'):
    """噬菌体协同变化页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果5：细菌-噬菌体协同变化" if lang == 'cn' else "Result 5: Bacteria-Phage Coordination"
    add_header(slide, prs, title)

    # 左侧内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    if lang == 'cn':
        bullets = [
            "噬菌体多样性显著降低 (P=0.032)",
            "",
            "细菌-噬菌体高度协同：",
            "  • Mantel r=0.52, P=0.001",
            "  • Procrustes P=0.002",
            "",
            "三角网络：17,848条边",
            "",
            "级联效应：",
            "  CagA → 细菌 → 噬菌体"
        ]
    else:
        bullets = [
            "Phage diversity significantly decreased (P=0.032)",
            "",
            "Bacteria-phage highly coordinated:",
            "  • Mantel r=0.52, P=0.001",
            "  • Procrustes P=0.002",
            "",
            "Tripartite network: 17,848 edges",
            "",
            "Cascade effect:",
            "  CagA → Bacteria → Phage"
        ]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        if "P=0.001" in bullet or "P=0.002" in bullet or "P=0.032" in bullet:
            p.font.bold = True
            p.font.color.rgb = COLORS['green']
        elif "级联" in bullet or "Cascade" in bullet:
            p.font.bold = True
            p.font.color.rgb = COLORS['primary_blue']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 右侧：级联效应示意图
    cascade_title = slide.shapes.add_textbox(Inches(5.5), Inches(1.4), Inches(4), Inches(0.4))
    tf = cascade_title.text_frame
    p = tf.paragraphs[0]
    p.text = "级联效应示意" if lang == 'cn' else "Cascade Effect"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 三个框：CagA → 细菌 → 噬菌体
    add_box_with_text(slide, 6.2, 1.9, 1.5, 0.6, "CagA\n感染" if lang == 'cn' else "CagA\nInfection", COLORS['red'], COLORS['white'], 12, True)

    arr1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.85), Inches(2.55), Inches(0.2), Inches(0.3))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = COLORS['dark_gray']
    arr1.line.fill.background()

    add_box_with_text(slide, 6.2, 2.9, 1.5, 0.6, "细菌\n重塑" if lang == 'cn' else "Bacteria\nRestructured", COLORS['primary_blue'], COLORS['white'], 12, True)

    arr2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.85), Inches(3.55), Inches(0.2), Inches(0.3))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = COLORS['dark_gray']
    arr2.line.fill.background()

    add_box_with_text(slide, 6.2, 3.9, 1.5, 0.6, "噬菌体\n协同变化" if lang == 'cn' else "Phage\nCoordinated", COLORS['purple'], COLORS['white'], 12, True)

    # 右侧补充说明
    note_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.6), Inches(4), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "CagA不直接影响病毒，通过细菌间接影响"
    else:
        p.text = "CagA affects phage indirectly via bacteria"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = COLORS['dark_gray']

    if lang == 'cn':
        notes = """细菌-噬菌体协同变化。（30秒）

虽然CagA不直接影响病毒组，但细菌的变化带动了噬菌体的改变。噬菌体多样性显著降低（P=0.032），细菌和噬菌体群落高度协同（Mantel P=0.001）。

这说明CagA的效应通过级联传递：CagA改变细菌，细菌改变噬菌体。"""
    else:
        notes = """Bacteria-phage coordination. (30s)

Although CagA doesn't directly affect viruses, bacterial changes drive phage changes. Phage diversity decreased (P=0.032), bacteria-phage highly coordinated (Mantel P=0.001).

Cascade effect: CagA → bacteria → phage."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_conclusion_slide(prs, title, conclusions, notes=""):
    """添加结论幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    for i, (num, text) in enumerate(conclusions):
        y_pos = 1.35 + i * 0.75

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['green']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.03), Inches(0.45), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_pos), Inches(8.4), Inches(0.65))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['dark_gray']

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_thanks_slide(prs, text, subtext=""):
    """添加致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_blue']
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_chinese_ppt(output_path, image_dir):
    """生成中文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. 封面
    add_title_slide(
        prs,
        "CagA依赖性肠道微生物组重塑\n及其功能冗余特征",
        "PC047 vCagAepitope 微生物组分析进展",
        "2026-01-26",
        "汇报人：龚宇航"
    )

    # 2. 背景1：CagA与肠道肿瘤
    add_background_slide_1(prs, 'cn')

    # 3. 背景2：Apc突变与G×E
    add_background_slide_2(prs, 'cn')

    # 4. 核心假说图解
    add_hypothesis_slide(prs, 'cn')

    # 5. 分析流程图
    add_analysis_pipeline_slide(prs, 'cn')

    # 6. 结果1：CagA显著重塑肠道菌群
    caga_effect_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "52_part5_caga_effect_by_genotype.png")
    add_content_slide(
        prs,
        "结果1：CagA显著重塑肠道菌群",
        [
            "• Beta多样性：**P=0.012**, R²=33.6%",
            "• Alpha多样性：P=0.73 (无变化)",
            "• → **选择性重组**，非全面破坏",
            "• 29个物种显著关联",
        ],
        notes="""CagA对肠道菌群的核心效应。（30秒）

CagA显著改变菌群结构（Beta P=0.012），但不改变整体多样性（Alpha P=0.73）。这说明CagA是选择性调控特定菌群，而非全面破坏。29个物种与CagA感染显著关联。""",
        image_path=caga_effect_path if os.path.exists(caga_effect_path) else None
    )

    # 7. 结果2：细菌vs病毒
    add_bacteria_vs_virus_slide(prs, 'cn')

    # 8. 结果3：功能冗余全层级验证（新增GO/PFAM）
    add_functional_redundancy_slide(prs, 'cn')

    # 9. 结果4：共现网络分析（新增）
    add_network_analysis_slide(prs, 'cn')

    # 10. 结果5：噬菌体协同变化
    add_phage_coordination_slide(prs, 'cn')

    # 11. 结论
    add_conclusion_slide(
        prs,
        "核心结论与后续计划",
        [
            (1, "CagA显著重塑肠道菌群：Beta多样性P=0.012，选择性调控"),
            (2, "功能冗余全层级验证：Pathway→KO→GO→PFAM，FDR均无显著"),
            (3, "网络模块结构瓦解：模块度降低62%，解释功能冗余机制"),
            (4, "Driver静默更替：益生菌↓致病菌↑，改变免疫刺激谱"),
            (5, "细菌-噬菌体协同：Mantel P=0.001，级联效应"),
            (6, "后续：代谢组学验证 + 整合肿瘤/T细胞表型数据"),
        ],
        notes="""核心结论。（45秒）

第一，CagA显著重塑菌群结构，选择性调控而非全面破坏。

第二，功能冗余得到全层级验证——从Pathway到PFAM，FDR均无显著差异。

第三，网络模块结构瓦解，模块度降低62%，这从网络层面解释了功能冗余机制。

第四，Driver物种静默更替，从益生菌转向致病菌。

第五，细菌和噬菌体高度协同，存在级联效应。

后续计划：代谢组学验证，整合肿瘤和T细胞表型数据。

感谢大家聆听！"""
    )

    # 12. 致谢
    add_thanks_slide(
        prs,
        "感谢聆听",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"中文版PPT已生成：{output_path}")


def generate_english_ppt(output_path, image_dir):
    """生成英文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title
    add_title_slide(
        prs,
        "CagA-Dependent Gut Microbiome Restructuring\nand Functional Redundancy",
        "PC047 vCagAepitope Microbiome Analysis Progress",
        "2026-01-26",
        "Presenter: Yuhang Gong"
    )

    # 2. Background 1
    add_background_slide_1(prs, 'en')

    # 3. Background 2
    add_background_slide_2(prs, 'en')

    # 4. Hypothesis
    add_hypothesis_slide(prs, 'en')

    # 5. Pipeline
    add_analysis_pipeline_slide(prs, 'en')

    # 6. Result 1: CagA Restructures Microbiota
    caga_effect_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "52_part5_caga_effect_by_genotype.png")
    add_content_slide(
        prs,
        "Result 1: CagA Restructures Gut Microbiota",
        [
            "• Beta diversity: **P=0.012**, R²=33.6%",
            "• Alpha diversity: P=0.73 (no change)",
            "• → **Selective restructuring**, not destruction",
            "• 29 species significantly associated",
        ],
        notes="""CagA core effect on gut microbiota. (30s)

CagA significantly alters community structure (Beta P=0.012), but does not change overall diversity (Alpha P=0.73). Selective regulation, not destruction. 29 species significantly associated.""",
        image_path=caga_effect_path if os.path.exists(caga_effect_path) else None
    )

    # 7. Result 2: Bacteria vs Virus
    add_bacteria_vs_virus_slide(prs, 'en')

    # 8. Result 3: Functional redundancy (with GO/PFAM)
    add_functional_redundancy_slide(prs, 'en')

    # 9. Result 4: Network analysis (new)
    add_network_analysis_slide(prs, 'en')

    # 10. Result 5: Phage coordination
    add_phage_coordination_slide(prs, 'en')

    # 11. Conclusions
    add_conclusion_slide(
        prs,
        "Conclusions & Future Directions",
        [
            (1, "CagA restructures gut microbiota: Beta P=0.012, selective regulation"),
            (2, "Functional redundancy at all levels: Pathway→KO→GO→PFAM, no FDR significance"),
            (3, "Network module collapse: Modularity -62%, explains redundancy mechanism"),
            (4, "Silent Driver shift: Probiotics↓ Pathobionts↑, altering immune stimulation"),
            (5, "Bacteria-phage coordination: Mantel P=0.001, cascade effect"),
            (6, "Next: Metabolomics validation + Tumor/T-cell phenotype integration"),
        ],
        notes="""Core conclusions. (45s)

First, CagA restructures microbiota through selective regulation.

Second, functional redundancy verified at all levels — Pathway to PFAM.

Third, network module collapse, modularity -62%, explaining functional redundancy.

Fourth, Driver species silently shifted from probiotics to pathobionts.

Fifth, bacteria-phage highly coordinated with cascade effect.

Next steps: metabolomics validation, tumor/T-cell phenotype integration.

Thank you for your attention!"""
    )

    # 12. Thanks
    add_thanks_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"英文版PPT已生成：{output_path}")


if __name__ == "__main__":
    base_dir = r"C:\Users\36094\Desktop\pc047oma"
    image_dir = os.path.join(base_dir, "analyses", "data")

    # 汇报日期和版本号
    meeting_date = "20260126"
    version = "v5"

    # 创建日期文件夹
    output_dir = os.path.join(base_dir, "PPT", meeting_date)
    os.makedirs(output_dir, exist_ok=True)

    # 输出路径
    chinese_path = os.path.join(output_dir, f"组会汇报_{version}.pptx")
    english_path = os.path.join(output_dir, f"GroupMeeting_{version}.pptx")

    generate_chinese_ppt(chinese_path, image_dir)
    generate_english_ppt(english_path, image_dir)

    print("\n生成完成！")
    print(f"中文版：{chinese_path}")
    print(f"英文版：{english_path}")

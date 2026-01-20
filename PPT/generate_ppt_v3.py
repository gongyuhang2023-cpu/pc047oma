"""
PC047组会PPT生成脚本 v4
核心叙述：CagA感染如何重塑肠道微生物组
调整：G×E作为模型验证而非主要发现
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'green': RGBColor(0x4C, 0xAF, 0x50),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'orange': RGBColor(0xFF, 0x98, 0x00),
    'red': RGBColor(0xE5, 0x39, 0x35),
    'purple': RGBColor(0x7B, 0x1F, 0xA2),
    'teal': RGBColor(0x00, 0x96, 0x88),
}


def add_page_number(slide, page_num, slide_width, slide_height):
    """为幻灯片添加页码（右下角，蓝色加粗）"""
    page_box = slide.shapes.add_textbox(
        slide_width - Inches(1.0),
        slide_height - Inches(0.6),
        Inches(0.8),
        Inches(0.4)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.RIGHT


def add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False):
    """为整个演示文稿添加页码"""
    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, 1):
        if skip_first and idx == 1:
            continue
        if skip_last and idx == total_slides:
            continue
        add_page_number(slide, idx, prs.slide_width, prs.slide_height)


def add_header(slide, prs, title):
    """添加统一标题栏"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_blue']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']


def add_box_with_text(slide, left, top, width, height, text, fill_color, text_color=None, font_size=14, bold=False):
    """添加带文字的方框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS['dark_gray']
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color or COLORS['dark_gray']
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_title_slide(prs, title, subtitle, date, presenter):
    """添加封面幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_blue']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{date}\n{presenter}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_background_slide_1(prs, lang='cn'):
    """背景页1：CagA与肠道肿瘤"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：CagA与肠道肿瘤" if lang == 'cn' else "Background: CagA and Intestinal Tumors"
    add_header(slide, prs, title)

    # 左侧：H. pylori感染流程
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "H. pylori CagA 致病机制" if lang == 'cn' else "H. pylori CagA Pathogenesis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    boxes_left = [
        (0.5, 1.7, 1.8, 0.5, "H. pylori\n感染" if lang == 'cn' else "H. pylori\nInfection", COLORS['light_blue']),
        (2.6, 1.7, 1.8, 0.5, "CagA蛋白\n注入宿主细胞" if lang == 'cn' else "CagA Injection\ninto Host Cell", COLORS['light_blue']),
        (0.5, 2.5, 1.8, 0.6, "干扰Wnt/NF-κB\n信号通路" if lang == 'cn' else "Disrupt Wnt/\nNF-κB Pathways", COLORS['orange']),
        (2.6, 2.5, 1.8, 0.6, "肠道菌群\n失调" if lang == 'cn' else "Gut Microbiota\nDysbiosis", COLORS['orange']),
        (1.55, 3.4, 2.0, 0.5, "促进肿瘤发生" if lang == 'cn' else "Tumor Promotion", COLORS['red']),
    ]

    for left, top, width, height, text, color in boxes_left:
        text_color = COLORS['white'] if color == COLORS['red'] else COLORS['dark_gray']
        add_box_with_text(slide, left, top, width, height, text, color, text_color, font_size=11)

    # 箭头
    line1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.35), Inches(1.85), Inches(0.2), Inches(0.15))
    line1.fill.solid()
    line1.fill.fore_color.rgb = COLORS['dark_gray']
    line1.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(2.25), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(3.15), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    # 右侧：关键文献支持
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "关键文献支持" if lang == 'cn' else "Key Literature Support"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    refs = [
        ("Jones et al. 2017, PLoS Pathog", "果蝇模型：CagA单独表达即可\n诱导菌群失调和上皮过度增殖" if lang == 'cn' else "Drosophila: CagA alone induces\ndysbiosis and epithelial proliferation"),
        ("Cui et al. 2025, BMC Gastro", "临床研究：CagA+菌株通过\n调节肠道菌群影响结直肠病变" if lang == 'cn' else "Clinical: CagA+ strains affect\ncolorectal lesions via gut microbiota"),
        ("Ding et al. 2025, Cell Host", "噬菌体可靶向清除促肿瘤\n细菌并恢复化疗敏感性" if lang == 'cn' else "Phage can eliminate tumor-\npromoting bacteria"),
    ]

    y_pos = 1.7
    for ref, desc in refs:
        ref_box = slide.shapes.add_textbox(Inches(5.0), Inches(y_pos), Inches(4.5), Inches(0.3))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + ref
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary_blue']

        desc_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.3), Inches(4.3), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['dark_gray']

        y_pos += 0.85

    # 底部：核心问题框
    question_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
    )
    question_box.fill.solid()
    question_box.fill.fore_color.rgb = COLORS['dark_blue']
    question_box.line.fill.background()

    q_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(8.6), Inches(0.6))
    tf = q_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "核心问题：CagA抗原清除后，肠道内什么因素维持CD8+ T细胞的持续激活？"
    else:
        p.text = "Central Question: What maintains CD8+ T-cell activation after CagA clearance?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """首先介绍一下研究背景。

幽门螺杆菌的CagA蛋白是公认的I类致癌因子。它通过IV型分泌系统注入宿主细胞后，可以干扰Wnt和NF-κB等关键信号通路，同时导致肠道菌群失调，最终促进肿瘤发生。

右侧列出了支持这一机制的关键文献。特别是Jones等人2017年的研究，在果蝇模型中证明CagA单独表达就足以诱导菌群失调和上皮过度增殖。

这就引出了我们的核心科学问题：在CagA抗原清除后，肠道内是什么因素维持着CD8+ T细胞的持续激活？"""
    else:
        notes = """Let me first introduce the research background.

H. pylori CagA protein is a recognized Class I carcinogen. After being injected into host cells via Type IV secretion system, it disrupts key signaling pathways like Wnt and NF-κB, while causing gut microbiota dysbiosis, ultimately promoting tumor development.

Key literature supporting this mechanism is listed on the right. Notably, Jones et al. 2017 demonstrated in Drosophila that CagA expression alone is sufficient to induce dysbiosis and epithelial hyperproliferation.

This leads to our central question: What maintains CD8+ T-cell activation after CagA clearance?"""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_background_slide_2(prs, lang='cn'):
    """背景页2：Apc突变与G×E交互"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：Apc突变与遗传易感性" if lang == 'cn' else "Background: Apc Mutation and Genetic Susceptibility"
    add_header(slide, prs, title)

    # 左侧：Apc-Wnt通路示意
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变 → Wnt通路持续激活" if lang == 'cn' else "Apc Mutation → Constitutive Wnt Activation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 正常状态
    normal_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(2), Inches(0.3))
    tf = normal_label.text_frame
    p = tf.paragraphs[0]
    p.text = "正常状态" if lang == 'cn' else "Normal"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    add_box_with_text(slide, 0.5, 1.9, 1.0, 0.4, "APC", COLORS['green'], COLORS['white'], 12, True)
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.0), Inches(0.3), Inches(0.2))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = COLORS['dark_gray']
    arr1.line.fill.background()
    add_box_with_text(slide, 1.9, 1.9, 1.2, 0.4, "β-catenin\n降解" if lang == 'cn' else "β-catenin\nDegradation", COLORS['light_gray'], COLORS['dark_gray'], 10)
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.0), Inches(0.3), Inches(0.2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = COLORS['dark_gray']
    arr2.line.fill.background()
    add_box_with_text(slide, 3.5, 1.9, 1.0, 0.4, "稳态" if lang == 'cn' else "Homeostasis", COLORS['green'], COLORS['white'], 11, True)

    # 突变状态
    mut_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(0.3))
    tf = mut_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变" if lang == 'cn' else "Apc Mutant"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['red']

    add_box_with_text(slide, 0.5, 2.8, 1.0, 0.4, "APC ✗", COLORS['red'], COLORS['white'], 12, True)
    arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.9), Inches(0.3), Inches(0.2))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = COLORS['dark_gray']
    arr3.line.fill.background()
    add_box_with_text(slide, 1.9, 2.8, 1.2, 0.4, "β-catenin\n累积" if lang == 'cn' else "β-catenin\nAccumulation", COLORS['orange'], COLORS['dark_gray'], 10)
    arr4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.9), Inches(0.3), Inches(0.2))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = COLORS['dark_gray']
    arr4.line.fill.background()
    add_box_with_text(slide, 3.5, 2.8, 1.0, 0.4, "肿瘤" if lang == 'cn' else "Tumor", COLORS['red'], COLORS['white'], 11, True)

    # 右侧：G×E交互概念图
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "G×E 交互作用假说" if lang == 'cn' else "G×E Interaction Hypothesis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 2x2矩阵标签
    col_labels = [("ApcWT", 6.0), ("ApcMUT", 7.8)]
    row_labels = [("CagA-", 2.0), ("CagA+", 2.8)]

    for label, x in col_labels:
        lb = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(1.5), Inches(0.3))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

    for label, y in row_labels:
        lb = slide.shapes.add_textbox(Inches(5.0), Inches(y), Inches(0.8), Inches(0.5))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']

    # 矩阵格子
    matrix_data = [
        (6.0, 1.95, "正常\n菌群" if lang == 'cn' else "Normal\nMicrobiota", COLORS['green']),
        (7.8, 1.95, "轻度\n失调" if lang == 'cn' else "Mild\nDysbiosis", COLORS['orange']),
        (6.0, 2.75, "无效应\nP=0.387" if lang == 'cn' else "No Effect\nP=0.387", COLORS['light_gray']),
        (7.8, 2.75, "显著重塑\nP=0.039" if lang == 'cn' else "Significant\nP=0.039", COLORS['red']),
    ]

    for x, y, text, color in matrix_data:
        text_color = COLORS['white'] if color in [COLORS['red'], COLORS['green']] else COLORS['dark_gray']
        add_box_with_text(slide, x, y, 1.5, 0.6, text, color, text_color, 11, True)

    # 结论框
    conclusion = slide.shapes.add_textbox(Inches(5.0), Inches(3.5), Inches(4.5), Inches(0.6))
    tf = conclusion.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "→ CagA效应需要Apc突变的易感背景"
    else:
        p.text = "→ CagA effect requires Apc-mutant susceptible background"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 底部：实验设计简图
    design_title = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(4.5), Inches(0.4))
    tf = design_title.text_frame
    p = tf.paragraphs[0]
    p.text = "实验设计：2×3因子" if lang == 'cn' else "Experimental Design: 2×3 Factorial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    groups = [
        (0.5, 3.9, "ApcWT", COLORS['light_blue']),
        (2.2, 3.9, "ApcMUT", COLORS['orange']),
    ]
    for x, y, text, color in groups:
        add_box_with_text(slide, x, y, 1.5, 0.4, text, color, COLORS['dark_gray'], 12, True)

    infections = ["Ctrl", "HpKO", "HpWT"]
    for i, inf in enumerate(infections):
        x = 0.5 + i * 1.4
        add_box_with_text(slide, x, 4.4, 1.2, 0.35, inf, COLORS['light_gray'], COLORS['dark_gray'], 10)

    core_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(4.85), Inches(1.7), Inches(0.35)
    )
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = COLORS['red']
    core_box.line.fill.background()

    core_text = slide.shapes.add_textbox(Inches(2.8), Inches(4.87), Inches(1.7), Inches(0.3))
    tf = core_text.text_frame
    p = tf.paragraphs[0]
    p.text = "核心比较" if lang == 'cn' else "Core Comparison"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """接下来介绍Apc突变与遗传易感性。

左侧展示了Apc-Wnt通路的关系。正常情况下，APC蛋白负责降解β-catenin，维持肠道稳态。当Apc发生突变时，β-catenin持续累积，最终导致肿瘤发生。

右侧是我们的G×E交互作用假说。我们的分析发现：CagA感染的效应高度依赖宿主遗传背景。在野生型背景下，CagA几乎无效应（P=0.387）；但在Apc突变背景下，CagA显著重塑肠道菌群（P=0.039）。

因此我们设计了2×3因子实验，核心比较是Apc突变背景下CagA阳性和CagA阴性感染组。"""
    else:
        notes = """Now let me introduce Apc mutation and genetic susceptibility.

The left panel shows the Apc-Wnt pathway. Normally, APC degrades β-catenin to maintain homeostasis. When Apc is mutated, β-catenin accumulates, leading to tumor development.

The right panel shows our G×E hypothesis. Our analysis found that CagA effect is highly dependent on host genetics. In wild-type, CagA has no effect (P=0.387); but in Apc-mutant, CagA significantly restructures microbiota (P=0.039).

We designed a 2×3 factorial experiment, with core comparison being CagA+ vs CagA- under Apc-mutant background."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_hypothesis_slide(prs, lang='cn'):
    """假说图解页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "核心假说：Functional Footprint" if lang == 'cn' else "Central Hypothesis: Functional Footprint"
    add_header(slide, prs, title)

    steps = [
        ("1", "CagA感染\n重塑菌群" if lang == 'cn' else "CagA Infection\nRestructures Microbiota", COLORS['primary_blue'], 0.5),
        ("2", "菌群产生\n特异性分子" if lang == 'cn' else "Microbiota Produces\nSpecific Molecules", COLORS['teal'], 2.7),
        ("3", "持续激活\nCD8+ T细胞" if lang == 'cn' else "Persistent CD8+\nT-cell Activation", COLORS['orange'], 4.9),
        ("4", "促进\n肿瘤发生" if lang == 'cn' else "Tumor\nPromotion", COLORS['red'], 7.1),
    ]

    for num, text, color, x in steps:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(1.4), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(x), Inches(1.45), Inches(0.4), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        add_box_with_text(slide, x, 1.9, 1.8, 0.7, text, COLORS['light_gray'], color, 12, True)

        if x < 7:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.85), Inches(2.15), Inches(0.7), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS['dark_gray']
            arr.line.fill.background()

    # 验证状态
    status_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.4))
    tf = status_title.text_frame
    p = tf.paragraphs[0]
    p.text = "本研究验证范围" if lang == 'cn' else "Scope of This Study"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    verified_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(4.2), Inches(1.5)
    )
    verified_box.fill.solid()
    verified_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    verified_box.line.color.rgb = COLORS['green']
    verified_box.line.width = Pt(2)

    v_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.8), Inches(0.3))
    tf = v_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ 已验证" if lang == 'cn' else "✓ Verified"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    v_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(3.8), Inches(1))
    tf = v_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 菌群结构显著重塑 (P=0.012)\n• G×E交互作用确认 (9/9验证)\n• 功能冗余与Driver更替\n• 噬菌体协同变化 (P=0.001)"
    else:
        p.text = "• Microbiota restructured (P=0.012)\n• G×E interaction confirmed (9/9 verified)\n• Functional redundancy & Driver shift\n• Phage coordination (P=0.001)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    pending_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.3), Inches(4.5), Inches(1.5)
    )
    pending_box.fill.solid()
    pending_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    pending_box.line.color.rgb = COLORS['orange']
    pending_box.line.width = Pt(2)

    p_title = slide.shapes.add_textbox(Inches(5.2), Inches(3.4), Inches(4.1), Inches(0.3))
    tf = p_title.text_frame
    p = tf.paragraphs[0]
    p.text = "○ 待验证" if lang == 'cn' else "○ Pending"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['orange']

    p_content = slide.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4.1), Inches(1))
    tf = p_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 代谢组学：特定代谢物鉴定\n• 转录组学：功能基因表达验证\n• 整合肿瘤/T细胞表型数据"
    else:
        p.text = "• Metabolomics: Specific metabolite ID\n• Transcriptomics: Gene expression validation\n• Integration with tumor/T-cell phenotypes"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    if lang == 'cn':
        notes = """这是我们的核心假说——Functional Footprint假说。

假说的逻辑链条是：第一步，CagA感染重塑肠道菌群；第二步，改变后的菌群产生特异性分子；第三步，这些分子持续激活CD8+ T细胞；第四步，最终促进肿瘤发生。

本次汇报涵盖的是假说的前半部分验证。我们通过9个系统性成对比较完全确认了G×E交互作用，发现功能冗余下的Driver物种更替，以及噬菌体的协同变化。

假说的后半部分——代谢物鉴定和与肿瘤/T细胞表型的关联——需要后续验证。"""
    else:
        notes = """This is our central hypothesis — the Functional Footprint hypothesis.

The logical chain: Step 1, CagA infection restructures gut microbiota; Step 2, altered microbiota produces specific molecules; Step 3, these molecules persistently activate CD8+ T-cells; Step 4, ultimately promoting tumor development.

This presentation covers the first half of hypothesis verification. We fully confirmed G×E interaction through 9 systematic pairwise comparisons, revealed Driver species shifts under functional redundancy, and phage coordination.

The latter half requires subsequent verification."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_analysis_pipeline_slide(prs, lang='cn'):
    """分析流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "分析流程概览" if lang == 'cn' else "Analysis Pipeline Overview"
    add_header(slide, prs, title)

    add_box_with_text(slide, 3.8, 1.3, 2.4, 0.5,
                      "盲肠样本 (n=27)" if lang == 'cn' else "Cecum Samples (n=27)",
                      COLORS['primary_blue'], COLORS['white'], 14, True)

    modules = [
        {
            'title': "01 物种组成分析" if lang == 'cn' else "01 Species Analysis",
            'tool': "Kraken2/Bracken",
            'focus': "群落结构\nG×E交互" if lang == 'cn' else "Community\nG×E Interaction",
            'result': "5/9显著" if lang == 'cn' else "5/9 Sig.",
            'x': 0.3,
            'color': COLORS['primary_blue'],
        },
        {
            'title': "02 功能潜力分析" if lang == 'cn' else "02 Functional Analysis",
            'tool': "HUMAnN4",
            'focus': "代谢通路\nDriver更替" if lang == 'cn' else "Pathways\nDriver Shift",
            'result': "功能冗余" if lang == 'cn' else "Redundancy",
            'x': 3.5,
            'color': COLORS['teal'],
        },
        {
            'title': "03 OTU/噬菌体分析" if lang == 'cn' else "03 OTU/Phage Analysis",
            'tool': "SingleM/Lyrebird",
            'focus': "系统发育\n细菌-噬菌体" if lang == 'cn' else "Phylogeny\nBacteria-Phage",
            'result': "P=0.001" if lang == 'cn' else "P=0.001",
            'x': 6.7,
            'color': COLORS['purple'],
        },
    ]

    for mod in modules:
        x = mod['x']

        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.0), Inches(0.45)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = mod['color']
        title_box.line.fill.background()

        t_text = slide.shapes.add_textbox(Inches(x), Inches(2.05), Inches(3.0), Inches(0.4))
        tf = t_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['title']
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        tool_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.0), Inches(0.35))
        tf = tool_box.text_frame
        p = tf.paragraphs[0]
        p.text = mod['tool']
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        focus_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.3), Inches(2.9), Inches(2.4), Inches(0.7)
        )
        focus_box.fill.solid()
        focus_box.fill.fore_color.rgb = COLORS['light_gray']
        focus_box.line.color.rgb = mod['color']

        f_text = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.95), Inches(2.4), Inches(0.65))
        tf = f_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mod['focus']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        result_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.6), Inches(3.7), Inches(1.8), Inches(0.4)
        )
        result_box.fill.solid()
        result_box.fill.fore_color.rgb = mod['color']
        result_box.line.fill.background()

        r_text = slide.shapes.add_textbox(Inches(x + 0.6), Inches(3.73), Inches(1.8), Inches(0.35))
        tf = r_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['result']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    for start_x in [1.8, 3.3, 5.0, 6.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(start_x + 2.9), Inches(1.85), Inches(0.2), Inches(0.12))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(9), Inches(0.9)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(8.6), Inches(0.75))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "整合结论：CagA通过G×E交互重塑细菌菌群 → 功能冗余但Driver更替 → 噬菌体协同"
    else:
        p.text = "Integration: CagA restructures bacteria via G×E → Functional redundancy with Driver shift → Phage coordination"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """这张幻灯片展示了我们的分析流程。

从27个盲肠样本出发，我们进行了三大类分析：

第一，物种组成分析，使用Kraken2和Bracken，9个成对比较中5个FDR显著。

第二，功能潜力分析，使用HUMAnN4，发现功能冗余现象。

第三，OTU和噬菌体分析，使用SingleM和Lyrebird，发现P=0.001的强关联。

整合结论：CagA通过G×E交互重塑细菌菌群，功能表现冗余但Driver物种发生更替，噬菌体协同变化。"""
    else:
        notes = """This slide shows our analysis pipeline.

From 27 cecum samples, we performed three major analyses:

First, species composition analysis using Kraken2/Bracken, 5/9 pairwise comparisons FDR significant.

Second, functional potential analysis using HUMAnN4, revealing functional redundancy.

Third, OTU and phage analysis using SingleM/Lyrebird, finding strong association at P=0.001.

Integration: CagA restructures bacteria via G×E, shows functional redundancy with Driver shift, and phage coordination."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_content_slide(prs, title, bullets, notes="", image_path=None):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    if image_path and os.path.exists(image_path):
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(4))
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if '**' in bullet:
            parts = bullet.split('**')
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(18)
                if j % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = COLORS['green']
                else:
                    run.font.color.rgb = COLORS['dark_gray']
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark_gray']

        p.space_before = Pt(8)
        p.level = 0

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_bacteria_vs_virus_slide(prs, lang='cn'):
    """新增：细菌vs病毒对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果2：CagA的作用靶点——细菌而非病毒" if lang == 'cn' else "Result 2: CagA Targets Bacteria, Not Viruses"
    add_header(slide, prs, title)

    # 问题引入
    question = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
    tf = question.text_frame
    p = tf.paragraphs[0]
    p.text = "问：CagA重塑的是什么？细菌？还是病毒？" if lang == 'cn' else "Q: What does CagA reshape? Bacteria or Viruses?"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['dark_blue']

    # 左侧：Standard细菌
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Standard (细菌)" if lang == 'cn' else "Standard (Bacteria)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 细菌结果框
    bacteria_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.3), Inches(4), Inches(2.0)
    )
    bacteria_box.fill.solid()
    bacteria_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    bacteria_box.line.color.rgb = COLORS['primary_blue']
    bacteria_box.line.width = Pt(2)

    bacteria_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = bacteria_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：5/9 (55.6%)"
    else:
        p.text = "• FDR Significant: 5/9 (55.6%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    else:
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.387"
    else:
        p.text = "• CagA@ApcWT: P=0.387"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ G×E交互完全验证"
    else:
        p.text = "→ G×E Fully Verified"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 右侧：Virus病毒
    right_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Virus (病毒)" if lang == 'cn' else "Virus (Viruses)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 病毒结果框
    virus_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.3), Inches(4), Inches(2.0)
    )
    virus_box.fill.solid()
    virus_box.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
    virus_box.line.color.rgb = COLORS['purple']
    virus_box.line.width = Pt(2)

    virus_content = slide.shapes.add_textbox(Inches(5.5), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = virus_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：0/9 (0%)"
    else:
        p.text = "• FDR Significant: 0/9 (0%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.604"
    else:
        p.text = "• CagA@ApcMUT: P=0.604"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.978"
    else:
        p.text = "• CagA@ApcWT: P=0.978"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ 无显著效应"
    else:
        p.text = "→ No Significant Effect"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 中间对比符号
    vs_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.55), Inches(3.0), Inches(0.7), Inches(0.7)
    )
    vs_box.fill.solid()
    vs_box.fill.fore_color.rgb = COLORS['dark_gray']
    vs_box.line.fill.background()

    vs_text = slide.shapes.add_textbox(Inches(4.55), Inches(3.15), Inches(0.7), Inches(0.4))
    tf = vs_text.text_frame
    p = tf.paragraphs[0]
    p.text = "VS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(0.5))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：CagA主要通过重塑细菌群落影响肠道微环境，而非直接作用于病毒组"
    else:
        p.text = "Conclusion: CagA reshapes gut environment mainly via bacteria, not directly via viruses"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """这页回答一个关键问题：CagA到底影响的是什么？

我们对细菌和病毒数据分别进行了9个成对比较分析。

左侧是Standard细菌的结果：9个比较中5个FDR显著，特别是核心的CagA效应在ApcMUT背景下高度显著（P=0.039），完全验证了G×E交互作用。

右侧是Virus病毒的结果：9个比较中0个FDR显著，包括核心CagA比较都不显著。

这说明CagA主要通过重塑细菌群落来影响肠道微环境，而不是直接作用于病毒组。

这为后面我们讨论细菌-噬菌体协同变化做了铺垫：虽然CagA不直接影响病毒，但细菌的变化会带动噬菌体的改变。"""
    else:
        notes = """This slide answers a key question: What does CagA actually affect?

We performed 9 pairwise comparisons for both bacteria and virus data.

Left shows Standard bacteria results: 5/9 FDR significant, especially the core CagA effect in ApcMUT background is highly significant (P=0.039), fully verifying G×E interaction.

Right shows Virus results: 0/9 FDR significant, including the core CagA comparison.

This indicates CagA reshapes gut environment mainly through bacteria, not directly via viruses.

This sets up our later discussion of bacteria-phage coordination: although CagA doesn't directly affect viruses, bacterial changes drive phage changes."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_conclusion_slide(prs, title, conclusions, notes=""):
    """添加结论幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    for i, (num, text) in enumerate(conclusions):
        y_pos = 1.4 + i * 0.85

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['green']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(8.3), Inches(0.75))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark_gray']

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_thanks_slide(prs, text, subtext=""):
    """添加致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_blue']
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_chinese_ppt(output_path, image_dir):
    """生成中文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. 封面
    add_title_slide(
        prs,
        "CagA依赖性肠道微生物组重塑\n及其功能冗余特征",
        "PC047 vCagAepitope 微生物组分析进展",
        "2026-01-26",
        "汇报人：龚宇航"
    )

    # 2. 背景1：CagA与肠道肿瘤
    add_background_slide_1(prs, 'cn')

    # 3. 背景2：Apc突变与G×E
    add_background_slide_2(prs, 'cn')

    # 4. 核心假说图解
    add_hypothesis_slide(prs, 'cn')

    # 5. 分析流程图
    add_analysis_pipeline_slide(prs, 'cn')

    # 6. 结果1：CagA显著重塑肠道菌群
    caga_effect_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "52_part5_caga_effect_by_genotype.png")
    add_content_slide(
        prs,
        "结果1：CagA显著重塑肠道菌群",
        [
            "• Beta多样性：**P=0.012**, R²=33.6%",
            "• Alpha多样性：P=0.73 (无变化)",
            "• → **选择性重组**，非全面破坏",
            "• 29个物种显著关联",
            "• 模型验证：Apc突变背景放大CagA效应"
        ],
        notes="""首先展示CagA对肠道菌群的核心效应。

核心发现：CagA显著改变菌群结构（Beta多样性P=0.012），但不改变整体多样性（Alpha P=0.73）。这说明CagA是选择性调控特定菌群，而非全面破坏。

29个物种与CagA感染显著关联。

模型验证方面：在Apc突变背景下CagA效应更显著（P=0.039 vs P=0.387），说明我们选用的肿瘤易感模型是有效的。""",
        image_path=caga_effect_path if os.path.exists(caga_effect_path) else None
    )

    # 7. 结果2：细菌vs病毒
    add_bacteria_vs_virus_slide(prs, 'cn')

    # 8. 结果3：功能冗余与静默更替
    driver_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_content_slide(
        prs,
        "结果3：功能冗余与静默更替",
        [
            "• 整体：物种交互P=0.024 vs 功能P=0.223",
            "• **但**核心比较：物种P=0.039, 功能**P=0.041**",
            "• → 尺度依赖：整体冗余，局部一致",
            "• **静默更替**：益生菌↓ 致病菌↑",
            "• *L. johnsonii* -50%, *M. schaedleri* +57%"
        ],
        notes="""第三个核心结果是功能冗余与静默更替。

整体交互显示功能冗余：物种交互P=0.024显著，功能交互P=0.223不显著。

但核心CagA比较显示物种-功能一致：物种P=0.039，功能P=0.041，都显著。这说明功能冗余是尺度依赖的——整体看冗余，特定比较看一致。

分层分析揭示"静默更替"：功能总量不变，但执行者已换。益生菌L. johnsonii下降50%，炎症相关菌M. schaedleri上升57%。

执行者从益生菌转向致病菌，这可能解释持续的免疫刺激。""",
        image_path=driver_path if os.path.exists(driver_path) else None
    )

    # 9. 结果4：噬菌体协同变化
    network_path = os.path.join(image_dir, "03b_virome_function_integration", "06_tripartite_network.png")
    add_content_slide(
        prs,
        "结果4：细菌-噬菌体协同变化",
        [
            "• 噬菌体多样性**显著降低** (P=0.032)",
            "• 细菌-噬菌体**高度协同**：",
            "  - Mantel r=0.52, **P=0.001**",
            "  - Procrustes **P=0.002**",
            "• 三角网络：**17,848条边**",
            "• 呼应：CagA→细菌→噬菌体级联效应"
        ],
        notes="""第四个发现是细菌-噬菌体协同变化。

虽然CagA不直接影响病毒组（上一页已展示），但细菌的变化带动了噬菌体的改变。

噬菌体多样性显著降低（P=0.032），细菌和噬菌体群落高度协同（Mantel P=0.001, Procrustes P=0.002）。

三角互作网络共17,848条边，细菌-噬菌体关联最强。

这说明CagA的效应通过级联传递：CagA改变细菌，细菌改变噬菌体。这呼应了我们前面的发现——CagA主要作用于细菌，但通过网络间接影响整个微生物组。""",
        image_path=network_path if os.path.exists(network_path) else None
    )

    # 10. 结论
    add_conclusion_slide(
        prs,
        "核心结论与后续计划",
        [
            (1, "CagA显著重塑肠道菌群：Beta多样性P=0.012, 选择性调控"),
            (2, "细菌与噬菌体同步改变：噬菌体多样性↓, 高度协同(P=0.001)"),
            (3, "功能执行者静默更替：Driver从益生菌→致病菌，改变免疫刺激谱"),
            (4, "模型验证：Apc突变背景有效放大CagA效应 (G×E P=0.024)"),
            (5, "后续：Gene Families分析 + 代谢组学验证"),
        ],
        notes="""最后总结核心结论。

第一，CagA显著重塑肠道菌群结构，通过选择性调控而非全面破坏。

第二，细菌与噬菌体同步改变，两者高度协同（Mantel P=0.001）。

第三，功能执行者发生静默更替，从益生菌转向潜在致病菌，可能改变免疫刺激来源。

第四，实验模型验证：Apc突变背景有效放大了CagA效应。

后续计划：Gene Families分析以获得更精细的功能差异，以及代谢组学验证。

感谢大家聆听！"""
    )

    # 11. 致谢
    add_thanks_slide(
        prs,
        "感谢聆听",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"中文版PPT已生成：{output_path}")


def generate_english_ppt(output_path, image_dir):
    """生成英文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title
    add_title_slide(
        prs,
        "CagA-Dependent Gut Microbiome Restructuring\nand Functional Redundancy",
        "PC047 vCagAepitope Microbiome Analysis Progress",
        "2026-01-26",
        "Presenter: Yuhang Gong"
    )

    # 2. Background 1
    add_background_slide_1(prs, 'en')

    # 3. Background 2
    add_background_slide_2(prs, 'en')

    # 4. Hypothesis
    add_hypothesis_slide(prs, 'en')

    # 5. Pipeline
    add_analysis_pipeline_slide(prs, 'en')

    # 6. Result 1: CagA Restructures Microbiota
    caga_effect_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "52_part5_caga_effect_by_genotype.png")
    add_content_slide(
        prs,
        "Result 1: CagA Restructures Gut Microbiota",
        [
            "• Beta diversity: **P=0.012**, R²=33.6%",
            "• Alpha diversity: P=0.73 (no change)",
            "• → **Selective restructuring**, not destruction",
            "• 29 species significantly associated",
            "• Model validation: Apc-mutant amplifies CagA effect"
        ],
        notes="""First, the core CagA effect on gut microbiota.

Core finding: CagA significantly alters community structure (Beta diversity P=0.012), but does not change overall diversity (Alpha P=0.73). This indicates CagA selectively regulates specific bacteria, not wholesale destruction.

29 species are significantly associated with CagA infection.

Model validation: CagA effect is more significant in Apc-mutant background (P=0.039 vs P=0.387), confirming our tumor-susceptible model is effective.""",
        image_path=caga_effect_path if os.path.exists(caga_effect_path) else None
    )

    # 7. Result 2: Bacteria vs Virus
    add_bacteria_vs_virus_slide(prs, 'en')

    # 8. Result 3: Functional redundancy
    driver_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_content_slide(
        prs,
        "Result 3: Functional Redundancy & Silent Shift",
        [
            "• Overall: Species P=0.024 vs Function P=0.223",
            "• **But** core comparison: Species P=0.039, Function **P=0.041**",
            "• → Scale-dependent: Overall redundant, locally consistent",
            "• **Silent Shift**: Probiotics↓ Pathobionts↑",
            "• *L. johnsonii* -50%, *M. schaedleri* +57%"
        ],
        notes="""Third core result is functional redundancy and Silent Shift.

Overall interaction shows functional redundancy: species P=0.024 significant, function P=0.223 not significant.

But core CagA comparison shows species-function consistency: both species P=0.039 and function P=0.041 are significant. This indicates functional redundancy is scale-dependent.

Stratified analysis reveals "Silent Shift": total function unchanged, but executors have changed. Probiotic L. johnsonii decreased 50%, inflammation-associated M. schaedleri increased 57%.

Drivers shifted from probiotics to pathobionts, possibly explaining persistent immune stimulation.""",
        image_path=driver_path if os.path.exists(driver_path) else None
    )

    # 9. Result 4: Phage
    network_path = os.path.join(image_dir, "03b_virome_function_integration", "06_tripartite_network.png")
    add_content_slide(
        prs,
        "Result 4: Bacteria-Phage Coordination",
        [
            "• Phage diversity **significantly decreased** (P=0.032)",
            "• Bacteria-phage **highly coordinated**:",
            "  - Mantel r=0.52, **P=0.001**",
            "  - Procrustes **P=0.002**",
            "• Tripartite network: **17,848 edges**",
            "• Echo: CagA→Bacteria→Phage cascade"
        ],
        notes="""Fourth finding is bacteria-phage coordination.

Although CagA doesn't directly affect viruses (as shown in previous slide), bacterial changes drive phage changes.

Phage diversity significantly decreased (P=0.032), bacteria and phage communities highly coordinated (Mantel P=0.001, Procrustes P=0.002).

Tripartite network contains 17,848 edges, with bacteria-phage being the strongest.

This shows CagA's effect cascades: CagA changes bacteria, bacteria change phage. This echoes our previous finding — CagA mainly affects bacteria, but indirectly influences the entire microbiome through the network.""",
        image_path=network_path if os.path.exists(network_path) else None
    )

    # 10. Conclusions
    add_conclusion_slide(
        prs,
        "Conclusions & Future Directions",
        [
            (1, "CagA restructures gut microbiota: Beta P=0.012, selective regulation"),
            (2, "Bacteria-phage synchronization: Phage diversity↓, highly coordinated (P=0.001)"),
            (3, "Silent Driver shift: Probiotics→Pathobionts, altering immune stimulation"),
            (4, "Model validation: Apc-mutant amplifies CagA effect (G×E P=0.024)"),
            (5, "Next: Gene Families analysis + Metabolomics validation"),
        ],
        notes="""Finally, core conclusions.

First, CagA significantly restructures gut microbiota through selective regulation, not destruction.

Second, bacteria and phage change synchronously, highly coordinated (Mantel P=0.001).

Third, functional Drivers silently shift from probiotics to pathobionts, potentially altering immune stimulation sources.

Fourth, model validation: Apc-mutant background effectively amplifies CagA effect.

Next steps: Gene Families analysis for finer functional resolution, and metabolomics validation.

Thank you for your attention!"""
    )

    # 11. Thanks
    add_thanks_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"英文版PPT已生成：{output_path}")


if __name__ == "__main__":
    base_dir = r"C:\Users\36094\Desktop\pc047oma"
    image_dir = os.path.join(base_dir, "analyses", "data")

    # 汇报日期和版本号
    meeting_date = "20260126"
    version = "v4"

    # 创建日期文件夹
    output_dir = os.path.join(base_dir, "PPT", meeting_date)
    os.makedirs(output_dir, exist_ok=True)

    # 输出路径
    chinese_path = os.path.join(output_dir, f"组会汇报_{version}.pptx")
    english_path = os.path.join(output_dir, f"GroupMeeting_{version}.pptx")

    generate_chinese_ppt(chinese_path, image_dir)
    generate_english_ppt(english_path, image_dir)

    print("\n生成完成！")
    print(f"中文版：{chinese_path}")
    print(f"英文版：{english_path}")

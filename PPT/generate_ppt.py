"""
PC047组会PPT生成脚本
使用python-pptx生成中英文双语版本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),    # #1465C0
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),       # #0D47A1
    'green': RGBColor(0x4C, 0xAF, 0x50),           # #4CAF50
    'dark_gray': RGBColor(0x33, 0x33, 0x33),       # #333333
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
}

def add_title_slide(prs, title, subtitle, date, presenter):
    """添加封面幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 添加背景色块（顶部）
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_blue']
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    # 日期和汇报人
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{date}\n{presenter}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, notes="", image_path=None):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题栏背景
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_blue']
    header.line.fill.background()

    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # 内容区域
    if image_path and os.path.exists(image_path):
        # 有图片时，左文右图布局
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(4))
        # 添加图片
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
    else:
        # 无图片时，全宽布局
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 处理加粗标记
        if '**' in bullet:
            parts = bullet.split('**')
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(18)
                if j % 2 == 1:  # 奇数索引是加粗部分
                    run.font.bold = True
                    run.font.color.rgb = COLORS['green']
                else:
                    run.font.color.rgb = COLORS['dark_gray']
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark_gray']

        p.space_before = Pt(8)
        p.level = 0

    # 添加讲稿到备注
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_table_slide(prs, title, headers, rows, notes=""):
    """添加表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题栏
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLORS['primary_blue']
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # 创建表格
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * table_rows)).table

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['dark_blue']
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.size = Pt(14)

    # 设置数据行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark_gray']

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_conclusion_slide(prs, title, conclusions, notes=""):
    """添加结论幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_blue']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # 结论框
    for i, (num, text) in enumerate(conclusions):
        y_pos = 1.4 + i * 1.1

        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['green']
        circle.line.fill.background()

        # 编号文字
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 结论文字
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(8.3), Inches(0.9))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_thanks_slide(prs, text, subtext=""):
    """添加致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 全屏背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_blue']
    bg.line.fill.background()

    # 致谢文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_chinese_ppt(output_path, image_dir):
    """生成中文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # 1. 封面
    add_title_slide(
        prs,
        "CagA依赖性肠道微生物组重塑\n及其功能冗余特征",
        "PC047 vCagAepitope 微生物组分析进展",
        "2026-01-17",
        "汇报人：龚宇航"
    )

    # 2. 研究背景与假说
    add_content_slide(
        prs,
        "研究背景与核心假说",
        [
            "• 核心科学问题：CagA抗原清除后，什么维持**CD8+ T细胞持续激活**？",
            "• **Functional Footprint假说**：被CagA永久改变的菌群产生的分子驱动T细胞激活",
            "• 实验模型：Apc突变小鼠 × H. pylori CagA+/- 感染",
            "• 核心比较：ApcMUT_HpWT (n=5) vs ApcMUT_HpKO (n=4)"
        ],
        notes="""大家好，今天我来汇报PC047项目的微生物组分析进展。

我们的核心科学问题是：在H. pylori感染清除后，肠道内是什么因素维持着CagA特异性CD8+ T细胞的持续激活？

基于此，我们提出了"Functional Footprint"假说——被CagA永久性改变的微生物群落可能产生某些非肽类分子或代谢物，驱动T细胞的持续激活。

为验证这一假说，我们使用Apc突变小鼠模型，设置了CagA阳性感染组和CagA阴性对照组的比较。"""
    )

    # 3. 实验设计
    add_table_slide(
        prs,
        "实验设计：2×3因子设计",
        ["组别", "基因型", "感染状态", "样本数", "说明"],
        [
            ["ApcWT_Ctrl", "野生型", "未感染", "5", "基线对照"],
            ["ApcWT_HpKO", "野生型", "CagA-", "5", "CagA阴性感染"],
            ["ApcWT_HpWT", "野生型", "CagA+", "5", "CagA阳性感染"],
            ["ApcMUT_Ctrl", "Apc突变", "未感染", "4", "排除ca08"],
            ["ApcMUT_HpKO", "Apc突变", "CagA-", "4", "核心对照组"],
            ["ApcMUT_HpWT", "Apc突变", "CagA+", "5", "核心实验组"],
        ],
        notes="""这是我们的实验分组设计。采用2×3因子设计，两个因素分别是Genotype（Apc野生型或突变型）和Infection（未感染、HpKO感染、HpWT感染）。

核心比较是ApcMUT背景下的HpWT和HpKO组——也就是在Apc突变这个肿瘤易感背景下，比较CagA阳性感染和CagA阴性感染的差异。

值得注意的是，我们的样本ca08因E. coli异常高丰度被剔除。"""
    )

    # 4. 物种组成分析
    pcoa_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "43_part4_pcoa_full_6groups.png")
    add_content_slide(
        prs,
        "物种组成分析：G×E交互作用",
        [
            "• 核心比较：Beta多样性 **P=0.012**，CagA显著重塑菌群结构",
            "• 2×3因子分析：Genotype×Infection **交互作用显著 (P=0.024)**",
            "• **关键发现**：CagA效应依赖Apc突变遗传背景",
            "    - ApcMUT背景：**P=0.007**, R²=33.6% ✓",
            "    - ApcWT背景：P=0.338，效应消失 ✗",
            "• 这是典型的**遗传-环境交互作用(G×E interaction)**"
        ],
        notes="""首先是物种组成分析的结果。

核心比较显示，CagA感染显著重塑了肠道菌群结构，Beta多样性PERMANOVA检验P值为0.012，解释了33%的群落变异。

更重要的是，当我们扩展到全部6组进行2×3因子分析时，发现了显著的Genotype与Infection交互作用，P值为0.024。

这意味着什么呢？进一步的分层分析揭示了关键发现：CagA的菌群重塑效应仅在Apc突变背景下显著，P值0.007；而在野生型背景下效应完全消失，P值0.338。

这是典型的遗传-环境交互作用模式，表明CagA需要宿主的"易感遗传背景"才能发挥菌群重塑作用。""",
        image_path=pcoa_path if os.path.exists(pcoa_path) else None
    )

    # 5. 功能冗余发现
    driver_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_content_slide(
        prs,
        "功能分析：功能冗余与静默更替",
        [
            "• **功能冗余现象**：物种交互P=0.024 vs 功能交互P=0.223",
            "• 物种显著改变，但整体代谢功能保持稳定",
            "• **Part 9新发现：静默更替(Silent Shift)**",
            "    - 165个通路-物种组合贡献变化>15%",
            "    - **益生菌贡献↓**：Lactobacillus johnsonii (-50%)",
            "    - **致病菌贡献↑**：Mucispirillum schaedleri (+57%)",
            "• 功能总量不变，但**执行者已从益生菌转向致病菌**"
        ],
        notes="""接下来是功能分析的核心发现。

我们观察到了典型的功能冗余现象：物种水平的交互作用显著(P=0.024)，但功能水平的交互作用不显著(P=0.223)。这意味着物种组成显著改变，但整体代谢功能保持稳定。

然而，Part 9的分层功能分析揭示了功能冗余表象下的重要变化——我们称之为"静默更替"现象。

分析发现165个通路-物种组合的贡献发生了超过15%的变化。具体来说，益生菌如Lactobacillus johnsonii的贡献下降了50%，而与肠道炎症相关的Mucispirillum schaedleri贡献上升了57%。

这意味着：虽然功能总量看起来不变，但执行这些功能的物种已经从益生菌转向了潜在致病菌。这为我们理解T细胞持续激活提供了新的机制解释。""",
        image_path=driver_path if os.path.exists(driver_path) else None
    )

    # 6. 噬菌体协同变化
    network_path = os.path.join(image_dir, "03b_virome_function_integration", "06_tripartite_network.png")
    add_content_slide(
        prs,
        "噬菌体分析：细菌-噬菌体协同变化",
        [
            "• **噬菌体多样性显著降低**：Shannon指数 P=0.032",
            "• 细菌-噬菌体群落**高度协同变化**：",
            "    - Mantel检验 r=0.52, **P=0.001**",
            "    - Procrustes分析 **P=0.002**",
            "• **三角互作网络**：17,848条边",
            "    - 细菌-噬菌体：7,544条（最强关联）",
            "    - 噬菌体-KO：6,604条",
            "• CagA效应通过**细菌-噬菌体-功能三角网络**传递"
        ],
        notes="""SingleM和Lyrebird分析揭示了噬菌体组的重要变化。

首先，CagA感染显著降低了噬菌体的alpha多样性，P值为0.032。这是Bracken物种分析无法发现的新现象。

更重要的是，细菌和噬菌体群落存在高度协同变化模式。Mantel检验显示两者距离矩阵高度相关，r=0.52，P值为0.001；Procrustes分析也证实了这种强烈关联。

我们构建的三角互作网络共有17,848条边，其中细菌-噬菌体的关联最强，达到7,544条。这表明CagA的效应可能通过这个细菌-噬菌体-功能三角网络传递和放大。""",
        image_path=network_path if os.path.exists(network_path) else None
    )

    # 7. 结论与后续
    add_conclusion_slide(
        prs,
        "核心结论与后续计划",
        [
            (1, "G×E交互作用确认：CagA菌群重塑效应依赖Apc突变背景 (P=0.007 vs P=0.338)"),
            (2, "功能冗余与静默更替：物种变但功能稳定，执行者从益生菌→致病菌"),
            (3, "噬菌体协同变化：CagA降低噬菌体多样性，细菌-噬菌体高度关联 (P=0.001)"),
            (4, "后续计划：代谢组学验证 + 整合肿瘤/T细胞表型数据"),
        ],
        notes="""最后总结一下我们的核心结论。

第一，我们首次在微生物组层面验证了G×E交互作用，CagA的菌群重塑效应需要Apc突变这个遗传背景。

第二，发现了功能冗余现象以及其背后的"静默更替"——功能稳定但执行者已经改变。

第三，揭示了噬菌体组的协同变化，细菌和噬菌体高度关联。

后续我们计划进行代谢组学分析来验证"静默更替"的代谢后果，并整合肿瘤和T细胞表型数据，直接验证假说的后半部分。

感谢大家的聆听！"""
    )

    # 8. 致谢
    add_thanks_slide(
        prs,
        "感谢聆听",
        "Questions & Discussion"
    )

    prs.save(output_path)
    print(f"中文版PPT已生成：{output_path}")


def generate_english_ppt(output_path, image_dir):
    """生成英文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title slide
    add_title_slide(
        prs,
        "CagA-Dependent Gut Microbiome Restructuring\nand Functional Redundancy",
        "PC047 vCagAepitope Microbiome Analysis Progress",
        "2026-01-17",
        "Presenter: Yuhang Gong"
    )

    # 2. Background
    add_content_slide(
        prs,
        "Background & Central Hypothesis",
        [
            "• Central question: What maintains **persistent CD8+ T-cell activation** after CagA clearance?",
            "• **Functional Footprint hypothesis**: Molecules from CagA-altered microbiota drive T-cell activation",
            "• Model: Apc-mutant mice × H. pylori CagA+/- infection",
            "• Core comparison: ApcMUT_HpWT (n=5) vs ApcMUT_HpKO (n=4)"
        ],
        notes="""Hello everyone, today I'll present our progress on microbiome analysis for the PC047 project.

Our central scientific question is: After H. pylori infection clearance, what maintains the persistent activation of CagA-specific CD8+ T-cells in the gut?

Based on this, we proposed the "Functional Footprint" hypothesis — that the microbiota permanently altered by CagA may produce certain non-peptide molecules or metabolites that drive persistent T-cell activation.

To test this hypothesis, we used an Apc-mutant mouse model, comparing CagA-positive infection group with CagA-negative control group."""
    )

    # 3. Experimental design
    add_table_slide(
        prs,
        "Experimental Design: 2×3 Factorial",
        ["Group", "Genotype", "Infection", "n", "Note"],
        [
            ["ApcWT_Ctrl", "Wild-type", "None", "5", "Baseline"],
            ["ApcWT_HpKO", "Wild-type", "CagA-", "5", "CagA-negative"],
            ["ApcWT_HpWT", "Wild-type", "CagA+", "5", "CagA-positive"],
            ["ApcMUT_Ctrl", "Apc-mutant", "None", "4", "Excluded ca08"],
            ["ApcMUT_HpKO", "Apc-mutant", "CagA-", "4", "Core control"],
            ["ApcMUT_HpWT", "Apc-mutant", "CagA+", "5", "Core experimental"],
        ],
        notes="""This is our experimental design. We employed a 2×3 factorial design with two factors: Genotype (Apc wild-type or mutant) and Infection (uninfected, HpKO, or HpWT).

The core comparison focuses on HpWT vs HpKO groups under the ApcMUT background — comparing CagA-positive and CagA-negative infection in a tumor-susceptible genetic background.

Note that sample ca08 was excluded due to abnormally high E. coli abundance."""
    )

    # 4. Species composition
    pcoa_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "43_part4_pcoa_full_6groups.png")
    add_content_slide(
        prs,
        "Species Analysis: G×E Interaction",
        [
            "• Core comparison: Beta diversity **P=0.012**, CagA restructures microbiota",
            "• 2×3 factorial: Genotype×Infection **interaction significant (P=0.024)**",
            "• **Key finding**: CagA effect depends on Apc-mutant background",
            "    - ApcMUT: **P=0.007**, R²=33.6% ✓",
            "    - ApcWT: P=0.338, effect absent ✗",
            "• Classic **Gene-Environment (G×E) interaction** pattern"
        ],
        notes="""First, the species composition analysis results.

The core comparison shows that CagA infection significantly restructures gut microbiota, with PERMANOVA P-value of 0.012, explaining 33% of community variation.

More importantly, when we extended to all 6 groups for 2×3 factorial analysis, we found significant Genotype×Infection interaction, P=0.024.

What does this mean? Stratified analysis revealed the key finding: CagA's microbiota-restructuring effect is only significant in Apc-mutant background (P=0.007), but completely absent in wild-type background (P=0.338).

This is a classic Gene-Environment interaction pattern, indicating that CagA requires a "susceptible genetic background" to exert its microbiota-restructuring effect.""",
        image_path=pcoa_path if os.path.exists(pcoa_path) else None
    )

    # 5. Functional redundancy
    driver_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_content_slide(
        prs,
        "Functional Analysis: Redundancy & Silent Shift",
        [
            "• **Functional redundancy**: Species interaction P=0.024 vs Function P=0.223",
            "• Species changed significantly, but metabolic function remains stable",
            "• **Part 9 discovery: Silent Shift phenomenon**",
            "    - 165 pathway-species combinations changed >15%",
            "    - **Probiotics decreased**: Lactobacillus johnsonii (-50%)",
            "    - **Pathobionts increased**: Mucispirillum schaedleri (+57%)",
            "• Same function, but **drivers shifted from probiotics to pathobionts**"
        ],
        notes="""Next, the core findings from functional analysis.

We observed classic functional redundancy: significant interaction at species level (P=0.024), but not at functional level (P=0.223). This means species composition changed significantly, but overall metabolic function remained stable.

However, Part 9 stratified analysis revealed important changes beneath the functional redundancy — what we call the "Silent Shift" phenomenon.

Analysis identified 165 pathway-species combinations with >15% contribution changes. Specifically, probiotics like Lactobacillus johnsonii decreased their contribution by 50%, while inflammation-associated Mucispirillum schaedleri increased by 57%.

This means: although total function appears unchanged, the species executing these functions have shifted from probiotics to potential pathogens. This provides a new mechanistic explanation for understanding persistent T-cell activation.""",
        image_path=driver_path if os.path.exists(driver_path) else None
    )

    # 6. Phage analysis
    network_path = os.path.join(image_dir, "03b_virome_function_integration", "06_tripartite_network.png")
    add_content_slide(
        prs,
        "Phage Analysis: Bacteria-Phage Coordination",
        [
            "• **Phage diversity significantly decreased**: Shannon P=0.032",
            "• Bacteria-phage communities **highly coordinated**:",
            "    - Mantel test r=0.52, **P=0.001**",
            "    - Procrustes analysis **P=0.002**",
            "• **Tripartite network**: 17,848 edges",
            "    - Bacteria-phage: 7,544 (strongest)",
            "    - Phage-KO: 6,604",
            "• CagA effect transmitted through **bacteria-phage-function network**"
        ],
        notes="""SingleM and Lyrebird analysis revealed important changes in the phageome.

First, CagA infection significantly reduced phage alpha diversity, with P-value of 0.032. This is a novel finding that Bracken species analysis could not detect.

More importantly, bacteria and phage communities show highly coordinated changes. Mantel test shows the two distance matrices are highly correlated (r=0.52, P=0.001); Procrustes analysis also confirmed this strong association.

Our tripartite interaction network contains 17,848 edges, with bacteria-phage associations being the strongest at 7,544 edges. This suggests that CagA's effect may be transmitted and amplified through this bacteria-phage-function tripartite network.""",
        image_path=network_path if os.path.exists(network_path) else None
    )

    # 7. Conclusions
    add_conclusion_slide(
        prs,
        "Conclusions & Future Directions",
        [
            (1, "G×E interaction confirmed: CagA effect depends on Apc-mutant background (P=0.007 vs P=0.338)"),
            (2, "Functional redundancy & Silent Shift: Species changed but function stable; drivers shifted"),
            (3, "Phage coordination: CagA reduces phage diversity; bacteria-phage highly associated (P=0.001)"),
            (4, "Next steps: Metabolomics validation + Integration with tumor/T-cell phenotype data"),
        ],
        notes="""Finally, let me summarize our core conclusions.

First, we confirmed G×E interaction at the microbiome level for the first time — CagA's microbiota-restructuring effect requires the Apc-mutant genetic background.

Second, we discovered functional redundancy and the underlying "Silent Shift" — function stable but drivers have changed.

Third, we revealed coordinated changes in the phageome, with bacteria and phage highly associated.

Next, we plan to perform metabolomics analysis to verify the metabolic consequences of "Silent Shift", and integrate tumor and T-cell phenotype data to directly validate the latter part of our hypothesis.

Thank you for your attention!"""
    )

    # 8. Thanks
    add_thanks_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    prs.save(output_path)
    print(f"英文版PPT已生成：{output_path}")


if __name__ == "__main__":
    # 设置路径
    base_dir = r"C:\Users\36094\Desktop\pc047oma"
    output_dir = os.path.join(base_dir, "PPT")
    image_dir = os.path.join(base_dir, "analyses", "data")

    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)

    # 生成PPT
    date_str = datetime.now().strftime("%Y-%m-%d")

    chinese_path = os.path.join(output_dir, f"组会汇报_{date_str}.pptx")
    english_path = os.path.join(output_dir, f"GroupMeeting_{date_str}.pptx")

    generate_chinese_ppt(chinese_path, image_dir)
    generate_english_ppt(english_path, image_dir)

    print("\n生成完成！")
    print(f"中文版：{chinese_path}")
    print(f"英文版：{english_path}")

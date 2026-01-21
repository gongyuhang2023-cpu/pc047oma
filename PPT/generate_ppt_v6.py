"""
PC047组会PPT生成脚本 v6
核心叙述：CagA感染如何重塑肠道微生物组
v6改进：加入高质量分析图片
- 结果1：使用综合Summary图（PCoA+效应量+成对比较）
- 结果3：加入Driver Species Shift图
- 结果5：加入Bacteria-Phage Procrustes图
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0x14, 0x65, 0xC0),
    'dark_blue': RGBColor(0x0D, 0x47, 0xA1),
    'green': RGBColor(0x4C, 0xAF, 0x50),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'orange': RGBColor(0xFF, 0x98, 0x00),
    'red': RGBColor(0xE5, 0x39, 0x35),
    'purple': RGBColor(0x7B, 0x1F, 0xA2),
    'teal': RGBColor(0x00, 0x96, 0x88),
}


def add_page_number(slide, page_num, slide_width, slide_height):
    """为幻灯片添加页码（右下角，蓝色加粗）"""
    page_box = slide.shapes.add_textbox(
        slide_width - Inches(1.0),
        slide_height - Inches(0.6),
        Inches(0.8),
        Inches(0.4)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.RIGHT


def add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False):
    """为整个演示文稿添加页码"""
    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, 1):
        if skip_first and idx == 1:
            continue
        if skip_last and idx == total_slides:
            continue
        add_page_number(slide, idx, prs.slide_width, prs.slide_height)


def add_header(slide, prs, title):
    """添加统一标题栏"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_blue']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']


def add_box_with_text(slide, left, top, width, height, text, fill_color, text_color=None, font_size=14, bold=False):
    """添加带文字的方框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS['dark_gray']
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color or COLORS['dark_gray']
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_title_slide(prs, title, subtitle, date, presenter):
    """添加封面幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_blue']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{date}\n{presenter}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_background_slide_1(prs, lang='cn'):
    """背景页1：CagA与肠道肿瘤"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：CagA与肠道肿瘤" if lang == 'cn' else "Background: CagA and Intestinal Tumors"
    add_header(slide, prs, title)

    # 左侧：H. pylori感染流程
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "H. pylori CagA 致病机制" if lang == 'cn' else "H. pylori CagA Pathogenesis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    boxes_left = [
        (0.5, 1.7, 1.8, 0.5, "H. pylori\n感染" if lang == 'cn' else "H. pylori\nInfection", COLORS['light_blue']),
        (2.6, 1.7, 1.8, 0.5, "CagA蛋白\n注入宿主细胞" if lang == 'cn' else "CagA Injection\ninto Host Cell", COLORS['light_blue']),
        (0.5, 2.5, 1.8, 0.6, "干扰Wnt/NF-κB\n信号通路" if lang == 'cn' else "Disrupt Wnt/\nNF-κB Pathways", COLORS['orange']),
        (2.6, 2.5, 1.8, 0.6, "肠道菌群\n失调" if lang == 'cn' else "Gut Microbiota\nDysbiosis", COLORS['orange']),
        (1.55, 3.4, 2.0, 0.5, "促进肿瘤发生" if lang == 'cn' else "Tumor Promotion", COLORS['red']),
    ]

    for left, top, width, height, text, color in boxes_left:
        text_color = COLORS['white'] if color == COLORS['red'] else COLORS['dark_gray']
        add_box_with_text(slide, left, top, width, height, text, color, text_color, font_size=11)

    # 箭头
    line1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.35), Inches(1.85), Inches(0.2), Inches(0.15))
    line1.fill.solid()
    line1.fill.fore_color.rgb = COLORS['dark_gray']
    line1.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(2.25), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    for x in [1.4, 3.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(3.15), Inches(0.15), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    # 右侧：关键文献支持
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "关键文献支持" if lang == 'cn' else "Key Literature Support"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    refs = [
        ("Jones et al. 2017, PLoS Pathog", "果蝇模型：CagA单独表达即可\n诱导菌群失调和上皮过度增殖" if lang == 'cn' else "Drosophila: CagA alone induces\ndysbiosis and epithelial proliferation"),
        ("Cui et al. 2025, BMC Gastro", "临床研究：CagA+菌株通过\n调节肠道菌群影响结直肠病变" if lang == 'cn' else "Clinical: CagA+ strains affect\ncolorectal lesions via gut microbiota"),
        ("Ding et al. 2025, Cell Host", "噬菌体可靶向清除促肿瘤\n细菌并恢复化疗敏感性" if lang == 'cn' else "Phage can eliminate tumor-\npromoting bacteria"),
    ]

    y_pos = 1.7
    for ref, desc in refs:
        ref_box = slide.shapes.add_textbox(Inches(5.0), Inches(y_pos), Inches(4.5), Inches(0.3))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + ref
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary_blue']

        desc_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.3), Inches(4.3), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['dark_gray']

        y_pos += 0.85

    # 底部：核心问题框
    question_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
    )
    question_box.fill.solid()
    question_box.fill.fore_color.rgb = COLORS['dark_blue']
    question_box.line.fill.background()

    q_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(8.6), Inches(0.6))
    tf = q_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "核心问题：CagA抗原清除后，肠道内什么因素维持CD8+ T细胞的持续激活？"
    else:
        p.text = "Central Question: What maintains CD8+ T-cell activation after CagA clearance?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """大家好，首先介绍一下研究背景。

我们都知道幽门螺杆菌是胃癌的一类致癌因子，而CagA蛋白是它最重要的毒力因子。

如左图所示，H. pylori感染后，CagA蛋白通过IV型分泌系统注入宿主细胞。进入细胞后，CagA会干扰两条关键信号通路：Wnt通路和NF-κB通路。与此同时，CagA还会导致肠道菌群失调。这两条路径最终都指向同一个结果——促进肿瘤发生。

右侧是支持这一机制的关键文献。Jones等人2017年在PLoS Pathogens发表的果蝇模型研究表明，CagA单独表达就足以诱导菌群失调和上皮过度增殖。

现在来看我们的核心科学问题：当CagA抗原被清除后，肠道内是什么因素在维持着CD8+ T细胞的持续激活？这正是本研究要回答的问题。"""
    else:
        notes = """Hello everyone, let me first introduce the research background.

We all know H. pylori is a class 1 carcinogen for gastric cancer, and CagA is its most important virulence factor.

As shown on the left, after H. pylori infection, CagA protein is injected into host cells via the type IV secretion system. Once inside, CagA disrupts two key signaling pathways: Wnt and NF-κB. Meanwhile, CagA also causes gut microbiota dysbiosis. Both pathways ultimately lead to tumor promotion.

On the right are key supporting studies. Jones et al. 2017 in PLoS Pathogens showed that CagA expression alone is sufficient to induce dysbiosis and epithelial hyperproliferation in Drosophila.

Now our central scientific question: After CagA antigen clearance, what maintains the persistent activation of CD8+ T cells in the gut? This is exactly what our study aims to answer."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_background_slide_2(prs, lang='cn'):
    """背景页2：Apc突变与G×E交互"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "研究背景：Apc突变与遗传易感性" if lang == 'cn' else "Background: Apc Mutation and Genetic Susceptibility"
    add_header(slide, prs, title)

    # 左侧：Apc-Wnt通路示意
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变 → Wnt通路持续激活" if lang == 'cn' else "Apc Mutation → Constitutive Wnt Activation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 正常状态
    normal_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(2), Inches(0.3))
    tf = normal_label.text_frame
    p = tf.paragraphs[0]
    p.text = "正常状态" if lang == 'cn' else "Normal"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    add_box_with_text(slide, 0.5, 1.9, 1.0, 0.4, "APC", COLORS['green'], COLORS['white'], 12, True)
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.0), Inches(0.3), Inches(0.2))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = COLORS['dark_gray']
    arr1.line.fill.background()
    add_box_with_text(slide, 1.9, 1.9, 1.2, 0.4, "β-catenin\n降解" if lang == 'cn' else "β-catenin\nDegradation", COLORS['light_gray'], COLORS['dark_gray'], 10)
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.0), Inches(0.3), Inches(0.2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = COLORS['dark_gray']
    arr2.line.fill.background()
    add_box_with_text(slide, 3.5, 1.9, 1.0, 0.4, "稳态" if lang == 'cn' else "Homeostasis", COLORS['green'], COLORS['white'], 11, True)

    # 突变状态
    mut_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(0.3))
    tf = mut_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Apc突变" if lang == 'cn' else "Apc Mutant"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['red']

    add_box_with_text(slide, 0.5, 2.8, 1.0, 0.4, "APC ✗", COLORS['red'], COLORS['white'], 12, True)
    arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.55), Inches(2.9), Inches(0.3), Inches(0.2))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = COLORS['dark_gray']
    arr3.line.fill.background()
    add_box_with_text(slide, 1.9, 2.8, 1.2, 0.4, "β-catenin\n累积" if lang == 'cn' else "β-catenin\nAccumulation", COLORS['orange'], COLORS['dark_gray'], 10)
    arr4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(2.9), Inches(0.3), Inches(0.2))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = COLORS['dark_gray']
    arr4.line.fill.background()
    add_box_with_text(slide, 3.5, 2.8, 1.0, 0.4, "肿瘤" if lang == 'cn' else "Tumor", COLORS['red'], COLORS['white'], 11, True)

    # 右侧：G×E交互概念图
    right_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "G×E 交互作用假说" if lang == 'cn' else "G×E Interaction Hypothesis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 2x2矩阵标签
    col_labels = [("ApcWT", 6.0), ("ApcMUT", 7.8)]
    row_labels = [("CagA-", 2.0), ("CagA+", 2.8)]

    for label, x in col_labels:
        lb = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(1.5), Inches(0.3))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

    for label, y in row_labels:
        lb = slide.shapes.add_textbox(Inches(5.0), Inches(y), Inches(0.8), Inches(0.5))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_gray']

    # 矩阵格子
    matrix_data = [
        (6.0, 1.95, "正常\n菌群" if lang == 'cn' else "Normal\nMicrobiota", COLORS['green']),
        (7.8, 1.95, "轻度\n失调" if lang == 'cn' else "Mild\nDysbiosis", COLORS['orange']),
        (6.0, 2.75, "无效应\nP=0.387" if lang == 'cn' else "No Effect\nP=0.387", COLORS['light_gray']),
        (7.8, 2.75, "显著重塑\nP=0.039" if lang == 'cn' else "Significant\nP=0.039", COLORS['red']),
    ]

    for x, y, text, color in matrix_data:
        text_color = COLORS['white'] if color in [COLORS['red'], COLORS['green']] else COLORS['dark_gray']
        add_box_with_text(slide, x, y, 1.5, 0.6, text, color, text_color, 11, True)

    # 结论框
    conclusion = slide.shapes.add_textbox(Inches(5.0), Inches(3.5), Inches(4.5), Inches(0.6))
    tf = conclusion.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "→ CagA效应需要Apc突变的易感背景"
    else:
        p.text = "→ CagA effect requires Apc-mutant susceptible background"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 底部：实验设计简图
    design_title = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(4.5), Inches(0.4))
    tf = design_title.text_frame
    p = tf.paragraphs[0]
    p.text = "实验设计：2×3因子" if lang == 'cn' else "Experimental Design: 2×3 Factorial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    groups = [
        (0.5, 3.9, "ApcWT", COLORS['light_blue']),
        (2.2, 3.9, "ApcMUT", COLORS['orange']),
    ]
    for x, y, text, color in groups:
        add_box_with_text(slide, x, y, 1.5, 0.4, text, color, COLORS['dark_gray'], 12, True)

    infections = ["Ctrl", "HpKO", "HpWT"]
    for i, inf in enumerate(infections):
        x = 0.5 + i * 1.4
        add_box_with_text(slide, x, 4.4, 1.2, 0.35, inf, COLORS['light_gray'], COLORS['dark_gray'], 10)

    core_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(4.85), Inches(1.7), Inches(0.35)
    )
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = COLORS['red']
    core_box.line.fill.background()

    core_text = slide.shapes.add_textbox(Inches(2.8), Inches(4.87), Inches(1.7), Inches(0.3))
    tf = core_text.text_frame
    p = tf.paragraphs[0]
    p.text = "核心比较" if lang == 'cn' else "Core Comparison"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """接下来介绍我们的实验模型设计。

左侧展示了Apc-Wnt通路的关系。正常情况下，APC蛋白负责降解β-catenin，保持Wnt信号稳定。但当Apc发生突变后，β-catenin无法被降解，导致Wnt信号持续激活，最终引发肠道肿瘤。这就是为什么Apc突变小鼠是研究肠道肿瘤的经典模型。

右侧是我们发现的一个关键现象：CagA效应高度依赖宿主遗传背景，这是典型的G×E交互作用。在野生型小鼠中，CagA几乎无效应，P值高达0.387。但在Apc突变小鼠中，CagA显著重塑肠道菌群，P值仅为0.039。这说明Apc突变提供了一个"易感背景"，放大了CagA的效应。

底部是我们的2×3因子实验设计：两种基因型乘以三种感染状态。红框标出的是核心比较——Apc突变背景下CagA阳性组和阴性组的对比。"""
    else:
        notes = """Now let me introduce our experimental model design.

The left side shows the Apc-Wnt pathway relationship. Normally, APC protein degrades β-catenin to maintain Wnt signaling stability. But when Apc mutates, β-catenin accumulates, leading to sustained Wnt activation and intestinal tumors. This is why Apc-mutant mice are the classic model for studying intestinal tumors.

On the right is a key finding: CagA effect is highly dependent on host genetics — a classic G×E interaction. In wild-type mice, CagA has almost no effect (P=0.387). But in Apc-mutant mice, CagA significantly restructures the microbiome (P=0.039). This suggests Apc mutation provides a "susceptible background" that amplifies CagA's effect.

At the bottom is our 2×3 factorial design: two genotypes times three infection states. The red box highlights our core comparison — CagA+ vs CagA- under Apc-mutant background."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_hypothesis_slide(prs, lang='cn'):
    """假说图解页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "核心假说：Functional Footprint" if lang == 'cn' else "Central Hypothesis: Functional Footprint"
    add_header(slide, prs, title)

    steps = [
        ("1", "CagA感染\n重塑菌群" if lang == 'cn' else "CagA Infection\nRestructures Microbiota", COLORS['primary_blue'], 0.5),
        ("2", "菌群产生\n特异性分子" if lang == 'cn' else "Microbiota Produces\nSpecific Molecules", COLORS['teal'], 2.7),
        ("3", "持续激活\nCD8+ T细胞" if lang == 'cn' else "Persistent CD8+\nT-cell Activation", COLORS['orange'], 4.9),
        ("4", "促进\n肿瘤发生" if lang == 'cn' else "Tumor\nPromotion", COLORS['red'], 7.1),
    ]

    for num, text, color, x in steps:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(1.4), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(x), Inches(1.45), Inches(0.4), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        add_box_with_text(slide, x, 1.9, 1.8, 0.7, text, COLORS['light_gray'], color, 12, True)

        if x < 7:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.85), Inches(2.15), Inches(0.7), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS['dark_gray']
            arr.line.fill.background()

    # 验证状态
    status_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.4))
    tf = status_title.text_frame
    p = tf.paragraphs[0]
    p.text = "本研究验证范围" if lang == 'cn' else "Scope of This Study"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    verified_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(4.2), Inches(1.5)
    )
    verified_box.fill.solid()
    verified_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    verified_box.line.color.rgb = COLORS['green']
    verified_box.line.width = Pt(2)

    v_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.8), Inches(0.3))
    tf = v_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ 已验证" if lang == 'cn' else "✓ Verified"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    v_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(3.8), Inches(1))
    tf = v_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 菌群结构显著重塑 (P=0.012)\n• G×E交互作用确认 (9/9验证)\n• 功能冗余 (GO/PFAM全验证)\n• 网络结构瓦解 (模块度-62%)"
    else:
        p.text = "• Microbiota restructured (P=0.012)\n• G×E interaction confirmed (9/9 verified)\n• Functional redundancy (GO/PFAM verified)\n• Network collapse (Modularity -62%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    pending_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.3), Inches(4.5), Inches(1.5)
    )
    pending_box.fill.solid()
    pending_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    pending_box.line.color.rgb = COLORS['orange']
    pending_box.line.width = Pt(2)

    p_title = slide.shapes.add_textbox(Inches(5.2), Inches(3.4), Inches(4.1), Inches(0.3))
    tf = p_title.text_frame
    p = tf.paragraphs[0]
    p.text = "○ 待验证" if lang == 'cn' else "○ Pending"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['orange']

    p_content = slide.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4.1), Inches(1))
    tf = p_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "• 代谢组学：特定代谢物鉴定\n• 转录组学：功能基因表达验证\n• 整合肿瘤/T细胞表型数据"
    else:
        p.text = "• Metabolomics: Specific metabolite ID\n• Transcriptomics: Gene expression validation\n• Integration with tumor/T-cell phenotypes"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    if lang == 'cn':
        notes = """这是我们的核心假说——Functional Footprint假说。

这个假说的逻辑链是这样的：首先，CagA感染会重塑肠道菌群的结构；然后，重塑后的菌群会产生一些特异性的分子或代谢物；这些分子会持续激活CD8+ T细胞；最终促进肿瘤发生。

绿框标出的是我们本次分析已经验证的部分。第一，菌群重塑已确认，Beta多样性P值仅0.012。第二，G×E交互作用验证成功，CagA效应确实依赖Apc遗传背景。第三，功能冗余得到全层级验证——从Pathway到KO到GO再到PFAM，FDR校正后都没有显著差异。第四，网络模块结构瓦解，模块度降低了62%。

橙框标出的是还需要后续验证的部分：代谢组学鉴定特定代谢物，以及整合肿瘤和T细胞表型数据。"""
    else:
        notes = """This is our central hypothesis — the Functional Footprint hypothesis.

The logical chain works like this: First, CagA infection restructures the gut microbiome; then, the restructured microbiome produces specific molecules or metabolites; these molecules persistently activate CD8+ T cells; ultimately promoting tumor development.

The green box shows what we've verified in this analysis. First, microbiome restructuring confirmed with Beta diversity P=0.012. Second, G×E interaction verified — CagA effect indeed depends on Apc genetic background. Third, functional redundancy verified at all levels — from Pathway to KO to GO to PFAM, no FDR-significant differences. Fourth, network module structure collapsed with modularity decreasing by 62%.

The orange box shows what still needs validation: metabolomics to identify specific metabolites, and integration with tumor and T-cell phenotype data."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_analysis_pipeline_slide(prs, lang='cn'):
    """分析流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "分析流程概览" if lang == 'cn' else "Analysis Pipeline Overview"
    add_header(slide, prs, title)

    add_box_with_text(slide, 3.8, 1.3, 2.4, 0.5,
                      "盲肠样本 (n=27)" if lang == 'cn' else "Cecum Samples (n=27)",
                      COLORS['primary_blue'], COLORS['white'], 14, True)

    modules = [
        {
            'title': "01 物种组成分析" if lang == 'cn' else "01 Species Analysis",
            'tool': "Kraken2/Bracken",
            'focus': "群落结构\nG×E交互" if lang == 'cn' else "Community\nG×E Interaction",
            'result': "5/9显著" if lang == 'cn' else "5/9 Sig.",
            'x': 0.3,
            'color': COLORS['primary_blue'],
        },
        {
            'title': "02 功能潜力分析" if lang == 'cn' else "02 Functional Analysis",
            'tool': "HUMAnN4",
            'focus': "Pathway→KO\n→GO→PFAM" if lang == 'cn' else "Pathway→KO\n→GO→PFAM",
            'result': "功能冗余" if lang == 'cn' else "Redundancy",
            'x': 3.5,
            'color': COLORS['teal'],
        },
        {
            'title': "03 网络/噬菌体" if lang == 'cn' else "03 Network/Phage",
            'tool': "igraph/Lyrebird",
            'focus': "共现网络\n细菌-噬菌体" if lang == 'cn' else "Co-occurrence\nBacteria-Phage",
            'result': "模块度-62%" if lang == 'cn' else "Mod -62%",
            'x': 6.7,
            'color': COLORS['purple'],
        },
    ]

    for mod in modules:
        x = mod['x']

        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.0), Inches(0.45)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = mod['color']
        title_box.line.fill.background()

        t_text = slide.shapes.add_textbox(Inches(x), Inches(2.05), Inches(3.0), Inches(0.4))
        tf = t_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['title']
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        tool_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.0), Inches(0.35))
        tf = tool_box.text_frame
        p = tf.paragraphs[0]
        p.text = mod['tool']
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        focus_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.3), Inches(2.9), Inches(2.4), Inches(0.7)
        )
        focus_box.fill.solid()
        focus_box.fill.fore_color.rgb = COLORS['light_gray']
        focus_box.line.color.rgb = mod['color']

        f_text = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.95), Inches(2.4), Inches(0.65))
        tf = f_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mod['focus']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER

        result_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.6), Inches(3.7), Inches(1.8), Inches(0.4)
        )
        result_box.fill.solid()
        result_box.fill.fore_color.rgb = mod['color']
        result_box.line.fill.background()

        r_text = slide.shapes.add_textbox(Inches(x + 0.6), Inches(3.73), Inches(1.8), Inches(0.35))
        tf = r_text.text_frame
        p = tf.paragraphs[0]
        p.text = mod['result']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    for start_x in [1.8, 3.3, 5.0, 6.5]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(start_x + 2.9), Inches(1.85), Inches(0.2), Inches(0.12))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS['dark_gray']
        arr.line.fill.background()

    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(9), Inches(0.9)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(8.6), Inches(0.75))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "整合结论：CagA重塑细菌菌群 → 功能冗余(全层级验证) → 网络结构瓦解 → 噬菌体协同"
    else:
        p.text = "Integration: CagA restructures bacteria → Functional redundancy (all levels) → Network collapse → Phage coordination"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """在介绍结果之前，先给大家概览一下我们的分析流程。

我们的分析分为三大模块。

第一个模块是物种组成分析。我们使用Kraken2和Bracken进行物种分类，包括细菌和病毒。同时使用SingleM和Lyrebird进行OTU水平的分析。核心发现：9个成对比较中有5个显著，证实了CagA的重塑效应。

第二个模块是功能潜力分析。我们使用HUMAnN4进行功能注释，然后用MaAsLin2进行差异分析。关键发现：从Pathway到KO到GO再到PFAM，四个功能层级都没有FDR显著差异，证实了功能冗余现象。

第三个模块是网络和噬菌体分析。我们构建了物种共现网络，发现模块度降低了62%。细菌-噬菌体Mantel检验显示两者高度协同，P值仅为0.001。

底部的整合结论是：CagA重塑细菌菌群，虽然存在功能冗余，但网络结构瓦解，噬菌体协同变化。接下来我会逐一展示这些发现。"""
    else:
        notes = """Before presenting results, let me give you an overview of our analysis pipeline.

Our analysis consists of three major modules.

The first module is species composition analysis. We used Kraken2 and Bracken for species classification, including bacteria and viruses, plus SingleM and Lyrebird for OTU-level analysis. Key finding: 5 out of 9 pairwise comparisons were significant, confirming CagA's restructuring effect.

The second module is functional potential analysis. We used HUMAnN4 for functional annotation, then MaAsLin2 for differential analysis. Key finding: from Pathway to KO to GO to PFAM, all four functional levels showed no FDR-significant differences, confirming functional redundancy.

The third module is network and phage analysis. We constructed species co-occurrence networks, finding modularity decreased by 62%. Bacteria-phage Mantel test showed high coordination with P=0.001.

The integrated conclusion at the bottom: CagA restructures bacteria, functional redundancy exists but network structure collapses, with phage coordination. I'll now walk through each finding."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_content_slide(prs, title, bullets, notes="", image_path=None):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    if image_path and os.path.exists(image_path):
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(4))
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if '**' in bullet:
            parts = bullet.split('**')
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(18)
                if j % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = COLORS['green']
                else:
                    run.font.color.rgb = COLORS['dark_gray']
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark_gray']

        p.space_before = Pt(8)
        p.level = 0

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_bacteria_vs_virus_slide(prs, lang='cn'):
    """新增：细菌vs病毒对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果2：CagA的作用靶点——细菌而非病毒" if lang == 'cn' else "Result 2: CagA Targets Bacteria, Not Viruses"
    add_header(slide, prs, title)

    # 问题引入
    question = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
    tf = question.text_frame
    p = tf.paragraphs[0]
    p.text = "问：CagA重塑的是什么？细菌？还是病毒？" if lang == 'cn' else "Q: What does CagA reshape? Bacteria or Viruses?"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['dark_blue']

    # 左侧：Standard细菌
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Standard (细菌)" if lang == 'cn' else "Standard (Bacteria)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 细菌结果框
    bacteria_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.3), Inches(4), Inches(2.0)
    )
    bacteria_box.fill.solid()
    bacteria_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    bacteria_box.line.color.rgb = COLORS['primary_blue']
    bacteria_box.line.width = Pt(2)

    bacteria_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = bacteria_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：5/9 (55.6%)"
    else:
        p.text = "• FDR Significant: 5/9 (55.6%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    else:
        p.text = "• CagA@ApcMUT: P=0.039 ⭐"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.387"
    else:
        p.text = "• CagA@ApcWT: P=0.387"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ G×E交互完全验证"
    else:
        p.text = "→ G×E Fully Verified"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # 右侧：Virus病毒
    right_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Virus (病毒)" if lang == 'cn' else "Virus (Viruses)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 病毒结果框
    virus_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.3), Inches(4), Inches(2.0)
    )
    virus_box.fill.solid()
    virus_box.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
    virus_box.line.color.rgb = COLORS['purple']
    virus_box.line.width = Pt(2)

    virus_content = slide.shapes.add_textbox(Inches(5.5), Inches(2.4), Inches(3.6), Inches(1.8))
    tf = virus_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "9个成对比较结果："
    else:
        p.text = "9 Pairwise Comparisons:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• FDR显著：0/9 (0%)"
    else:
        p.text = "• FDR Significant: 0/9 (0%)"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcMUT: P=0.604"
    else:
        p.text = "• CagA@ApcMUT: P=0.604"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "• CagA@ApcWT: P=0.978"
    else:
        p.text = "• CagA@ApcWT: P=0.978"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark_gray']

    p = tf.add_paragraph()
    if lang == 'cn':
        p.text = "→ 无显著效应"
    else:
        p.text = "→ No Significant Effect"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # 中间对比符号
    vs_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.55), Inches(3.0), Inches(0.7), Inches(0.7)
    )
    vs_box.fill.solid()
    vs_box.fill.fore_color.rgb = COLORS['dark_gray']
    vs_box.line.fill.background()

    vs_text = slide.shapes.add_textbox(Inches(4.55), Inches(3.15), Inches(0.7), Inches(0.4))
    tf = vs_text.text_frame
    p = tf.paragraphs[0]
    p.text = "VS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(0.5))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：CagA主要通过重塑细菌群落影响肠道微环境，而非直接作用于病毒组"
    else:
        p.text = "Conclusion: CagA reshapes gut environment mainly via bacteria, not directly via viruses"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """既然CagA能重塑肠道微生物组，那一个自然的问题是：它重塑的到底是什么？是细菌？还是病毒？

我们对细菌和病毒分别进行了9个成对比较的系统分析，结果非常清晰。

左侧是细菌的分析结果。9个比较中有5个是FDR校正后显著的，超过一半。最关键的核心比较——CagA效应在Apc突变背景下——P值仅0.039，完全验证了G×E交互作用。G×E交互也是显著的。

右侧是病毒的分析结果。9个比较中有0个是FDR显著的，一个都没有！最关键的核心CagA比较，P值高达0.604，完全不显著。G×E交互也不显著。

这个对比告诉我们什么？CagA主要通过重塑细菌群落来影响肠道微环境，而不是直接作用于病毒组。换句话说，细菌是CagA的直接靶点，病毒的变化可能是细菌变化的下游效应。这就引出了后面噬菌体协同变化的分析。"""
    else:
        notes = """Since CagA can reshape the gut microbiome, a natural question is: what exactly does it reshape? Bacteria? Or viruses?

We systematically performed 9 pairwise comparisons for both bacteria and viruses. The results are very clear.

On the left are bacteria results. Out of 9 comparisons, 5 are FDR-significant — over half. Most importantly, the core comparison — CagA effect under Apc-mutant background — has P=0.039, fully verifying the G×E interaction.

On the right are virus results. Out of 9 comparisons, 0 are FDR-significant — none at all! The core CagA comparison has P=0.604, completely non-significant. G×E interaction is also not significant.

What does this comparison tell us? CagA reshapes the gut environment mainly by restructuring bacteria, not by directly acting on viruses. In other words, bacteria are CagA's direct target, and viral changes may be downstream effects of bacterial changes. This leads to our later analysis of bacteria-phage coordination."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_functional_redundancy_slide(prs, lang='cn', image_path=None):
    """新增：功能冗余全层级验证页（含GO/PFAM），支持Driver图片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果3：功能冗余——全层级验证" if lang == 'cn' else "Result 3: Functional Redundancy — All Levels Verified"
    add_header(slide, prs, title)

    # 左侧：功能层级验证表
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HUMAnN4 功能层级验证" if lang == 'cn' else "HUMAnN4 Functional Hierarchy"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 表格数据
    table_data = [
        ("层级", "特征数", "P<0.05", "FDR<0.05", "结论"),
        ("Pathway", "~500", "多个", "0", "✓ 冗余"),
        ("KO", "~6,000", "多个", "0", "✓ 冗余"),
        ("GO", "~13,000", "1,577", "0", "✓ 冗余"),
        ("PFAM", "~7,600", "849", "0", "✓ 冗余"),
    ]

    # 创建表格
    table_top = 1.7
    row_height = 0.4
    col_widths = [1.2, 1.1, 0.9, 1.0, 0.9]

    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            x = 0.4 + sum(col_widths[:col_idx])
            y = table_top + row_idx * row_height

            if row_idx == 0:
                # 表头
                cell_bg = COLORS['primary_blue']
                text_color = COLORS['white']
                font_bold = True
            elif "✓" in str(cell_text):
                cell_bg = RGBColor(0xE8, 0xF5, 0xE9)
                text_color = COLORS['green']
                font_bold = True
            else:
                cell_bg = COLORS['light_gray'] if row_idx % 2 == 0 else COLORS['white']
                text_color = COLORS['dark_gray']
                font_bold = False

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[col_idx]), Inches(row_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = cell_bg
            cell.line.color.rgb = COLORS['dark_gray']
            cell.line.width = Pt(0.5)

            cell_text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(col_widths[col_idx]), Inches(0.3))
            tf = cell_text_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(10)
            p.font.bold = font_bold
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

    # 右侧：Driver更替 - 使用图片或文字
    right_title = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "静默更替：Driver物种改变" if lang == 'cn' else "Silent Shift: Driver Species Change"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # 如果有图片，显示图片；否则显示文字
    if image_path and os.path.exists(image_path):
        # 使用Driver Species Shift图片
        slide.shapes.add_picture(image_path, Inches(5.3), Inches(1.7), width=Inches(4.4))
    else:
        # Driver更替框（备用文字版）
        driver_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.7), Inches(4), Inches(2.3)
        )
        driver_box.fill.solid()
        driver_box.fill.fore_color.rgb = RGBColor(0xFB, 0xE9, 0xE7)
        driver_box.line.color.rgb = COLORS['orange']
        driver_box.line.width = Pt(2)

        driver_content = slide.shapes.add_textbox(Inches(5.7), Inches(1.8), Inches(3.6), Inches(2.1))
        tf = driver_content.text_frame
        tf.word_wrap = True

        if lang == 'cn':
            lines = [
                "功能总量不变，但执行者改变：",
                "",
                "益生菌 ↓",
                "  • L. johnsonii: -50%",
                "",
                "炎症相关菌 ↑",
                "  • M. schaedleri: +57%",
            ]
        else:
            lines = [
                "Total function unchanged, executors changed:",
                "",
                "Probiotics ↓",
                "  • L. johnsonii: -50%",
                "",
                "Inflammation-associated ↑",
                "  • M. schaedleri: +57%",
            ]

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            if "↓" in line and "益生菌" in line:
                p.font.color.rgb = COLORS['green']
                p.font.bold = True
            elif "↑" in line and ("炎症" in line or "Inflammation" in line):
                p.font.color.rgb = COLORS['red']
                p.font.bold = True
            else:
                p.font.color.rgb = COLORS['dark_gray']

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(1.0)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(8.6), Inches(0.8))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：从Pathway→KO→GO→PFAM全层级FDR校正后均无显著差异\n功能冗余得到全面验证，但功能执行者已静默更替"
    else:
        p.text = "Conclusion: No FDR-significant differences at any level (Pathway→KO→GO→PFAM)\nFunctional redundancy fully verified, but executors silently shifted"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """接下来是一个非常有意思的发现：功能冗余现象。

左侧的表格展示了我们对四个功能层级的系统验证。从最粗的Pathway层级，到KO基因功能，再到GO基因本体，最后到PFAM蛋白结构域——这是从粗到细的四个功能层面。

结果非常一致：每个层级都有很多P值小于0.05的特征，但经过FDR校正后，显著差异的数量都是零！这意味着什么？虽然物种组成发生了显著变化，但整体的功能潜力保持稳定。这就是生态学中的"功能冗余"现象——不同物种可以执行相同的功能。

但故事到这里还没有结束。右侧图片展示了我们的分层功能分析，也就是HUMAnN的stratified analysis。这张彩色堆叠图展示了关键通路中不同物种的贡献比例变化。

我们发现了一个重要现象，我称之为"静默更替"。虽然功能总量不变，但执行这些功能的物种已经换了！左侧的益生菌，比如Lactobacillus johnsonii，贡献下降了50%；右侧的炎症相关菌，比如Mucispirillum schaedleri，贡献上升了57%。

这意味着什么？虽然代谢功能看起来一样，但产生这些功能的"工人"已经从好菌变成了坏菌。这对免疫系统来说是完全不同的信号。"""
    else:
        notes = """Next is a very interesting finding: functional redundancy.

The table on the left shows our systematic verification across four functional levels. From the coarsest Pathway level, to KO gene functions, to GO gene ontology, to PFAM protein domains — these are four levels from coarse to fine.

The results are remarkably consistent: each level has many features with P<0.05, but after FDR correction, the number of significant differences is zero! What does this mean? Although species composition changed significantly, overall functional potential remains stable. This is the ecological phenomenon of "functional redundancy" — different species can perform the same functions.

But the story doesn't end here. The image on the right shows our stratified functional analysis from HUMAnN. This colorful stacked bar chart shows how different species' contributions to key pathways changed.

We discovered an important phenomenon I call "Silent Shift." Although total function remains unchanged, the species executing these functions have changed! On the left, probiotics like Lactobacillus johnsonii decreased their contribution by 50%; on the right, inflammation-associated bacteria like Mucispirillum schaedleri increased by 57%.

What does this mean? Although metabolic functions look the same, the "workers" producing these functions have changed from good bacteria to bad bacteria. For the immune system, this is a completely different signal."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_network_analysis_slide(prs, lang='cn'):
    """新增：共现网络分析页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果4：共现网络——模块结构瓦解" if lang == 'cn' else "Result 4: Co-occurrence Network — Module Collapse"
    add_header(slide, prs, title)

    # 左侧：网络拓扑对比
    left_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "网络拓扑变化 (Top 100物种)" if lang == 'cn' else "Network Topology Change (Top 100 Species)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # HpKO (对照) 框
    hpko_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.75), Inches(2.2), Inches(1.6)
    )
    hpko_box.fill.solid()
    hpko_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    hpko_box.line.color.rgb = COLORS['green']
    hpko_box.line.width = Pt(2)

    hpko_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(2), Inches(0.3))
    tf = hpko_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HpKO (对照)" if lang == 'cn' else "HpKO (Control)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['green']

    hpko_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(2), Inches(1.2))
    tf = hpko_content.text_frame
    tf.word_wrap = True
    if lang == 'cn':
        lines = ["节点: 80", "边数: 1,280", "模块度: 0.468", "→ 清晰模块结构"]
    else:
        lines = ["Nodes: 80", "Edges: 1,280", "Modularity: 0.468", "→ Clear modules"]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "模块度" in line or "Modularity" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['green']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.7), Inches(2.4), Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS['red']
    arrow.line.fill.background()

    # HpWT (CagA+) 框
    hpwt_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.75), Inches(2.2), Inches(1.6)
    )
    hpwt_box.fill.solid()
    hpwt_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    hpwt_box.line.color.rgb = COLORS['red']
    hpwt_box.line.width = Pt(2)

    hpwt_title = slide.shapes.add_textbox(Inches(3.3), Inches(1.8), Inches(2), Inches(0.3))
    tf = hpwt_title.text_frame
    p = tf.paragraphs[0]
    p.text = "HpWT (CagA+)" if lang == 'cn' else "HpWT (CagA+)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['red']

    hpwt_content = slide.shapes.add_textbox(Inches(3.3), Inches(2.1), Inches(2), Inches(1.2))
    tf = hpwt_content.text_frame
    tf.word_wrap = True
    if lang == 'cn':
        lines = ["节点: 100", "边数: 2,190 (+71%)", "模块度: 0.177", "→ 模块结构瓦解"]
    else:
        lines = ["Nodes: 100", "Edges: 2,190 (+71%)", "Modularity: 0.177", "→ Module collapse"]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "模块度" in line or "Modularity" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['red']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 模块度变化强调框
    mod_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.5), Inches(5.0), Inches(0.7)
    )
    mod_box.fill.solid()
    mod_box.fill.fore_color.rgb = COLORS['orange']
    mod_box.line.fill.background()

    mod_text = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(4.8), Inches(0.5))
    tf = mod_text.text_frame
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "模块度: 0.468 → 0.177 (降低62%)"
    else:
        p.text = "Modularity: 0.468 → 0.177 (-62%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 右侧：Hub物种
    right_title = slide.shapes.add_textbox(Inches(5.8), Inches(1.3), Inches(3.7), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Hub物种 (网络枢纽)" if lang == 'cn' else "Hub Species (Network Hubs)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    hub_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.7), Inches(3.7), Inches(2.5)
    )
    hub_box.fill.solid()
    hub_box.fill.fore_color.rgb = COLORS['light_gray']
    hub_box.line.color.rgb = COLORS['purple']
    hub_box.line.width = Pt(1)

    hub_content = slide.shapes.add_textbox(Inches(5.95), Inches(1.8), Inches(3.4), Inches(2.3))
    tf = hub_content.text_frame
    tf.word_wrap = True

    if lang == 'cn':
        hub_lines = [
            "Top 5 Hub物种（按度排序）：",
            "• Blautia obeum (度=44)",
            "• Roseburia hominis (度=42)",
            "• Hungatella hathewayi (度=42)",
            "",
            "特征：主要为厌氧梭菌类",
            "功能：SCFA生产者",
            "",
            "桥梁物种：",
            "• P. distasonis (介数=402)"
        ]
    else:
        hub_lines = [
            "Top 5 Hub Species (by degree):",
            "• Blautia obeum (d=44)",
            "• Roseburia hominis (d=42)",
            "• Hungatella hathewayi (d=42)",
            "",
            "Feature: Mainly anaerobic Clostridia",
            "Function: SCFA producers",
            "",
            "Bridge species:",
            "• P. distasonis (btw=402)"
        ]

    for i, line in enumerate(hub_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        if "Top" in line or "桥梁" in line or "Bridge" in line or "特征" in line or "Feature" in line:
            p.font.bold = True
            p.font.color.rgb = COLORS['purple']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 底部结论
    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.35), Inches(9), Inches(0.85)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = COLORS['dark_blue']
    conclusion_box.line.fill.background()

    c_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.45), Inches(8.6), Inches(0.7))
    tf = c_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if lang == 'cn':
        p.text = "结论：CagA导致网络模块度降低62%，物种间界限模糊\n→ 生态位重叠增加 → 解释功能冗余机制"
    else:
        p.text = "Conclusion: CagA causes 62% modularity drop, blurred species boundaries\n→ Increased niche overlap → Explains functional redundancy mechanism"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if lang == 'cn':
        notes = """这一页我们来看共现网络分析，它揭示了一个重要现象：网络模块结构的瓦解。

什么是模块度？简单来说，高模块度意味着物种各司其职，功能分化明显；低模块度意味着物种混在一起，功能界限模糊。

左侧表格对比了两组的网络特征。对照组HpKO：80个节点，1280条边，模块度0.468，这是一个有清晰功能模块的健康网络。CagA感染组HpWT：100个节点，2190条边——增加了71%！但最关键的是模块度，从0.468降到了0.177，降幅达62%。

这个62%的模块度下降意味着什么？物种间的界限变模糊了，生态位重叠增加了。原本各司其职的物种，现在功能开始重叠。这从网络层面完美解释了前面看到的功能冗余现象——为什么物种变了但功能稳定，因为物种之间的功能分工被打乱了。

右侧展示了网络中的Hub物种，也就是连接度最高的关键物种。有意思的是，这些Hub物种主要是厌氧梭菌类，比如Blautia、Roseburia等，它们都是短链脂肪酸的主要生产者。另外值得注意的是Parabacteroides distasonis，它的介数中心性特别高，是连接不同功能模块的"桥梁物种"，可能是CagA影响免疫调节的关键节点。"""
    else:
        notes = """This slide shows co-occurrence network analysis, revealing an important phenomenon: network module structure collapse.

What is modularity? Simply put, high modularity means species specialize and functions are differentiated; low modularity means species mix together with blurred functional boundaries.

The left table compares network characteristics between groups. Control HpKO: 80 nodes, 1,280 edges, modularity 0.468 — a healthy network with clear functional modules. CagA+ HpWT: 100 nodes, 2,190 edges — a 71% increase! But most critically, modularity dropped from 0.468 to 0.177, a 62% decrease.

What does this 62% modularity drop mean? Species boundaries became blurred, niche overlap increased. Species that used to specialize now have overlapping functions. This perfectly explains the functional redundancy we saw earlier at the network level — why species changed but functions remained stable, because functional division among species was disrupted.

The right side shows Hub species — the most connected key species in the network. Interestingly, these Hub species are mainly anaerobic Clostridia like Blautia and Roseburia — major short-chain fatty acid producers. Also noteworthy is Parabacteroides distasonis with exceptionally high betweenness centrality — a "bridge species" connecting different functional modules, potentially a key node for CagA's immune modulation effect."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_phage_coordination_slide(prs, lang='cn', image_path=None):
    """噬菌体协同变化页，支持Procrustes图片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title = "结果5：细菌-噬菌体协同变化" if lang == 'cn' else "Result 5: Bacteria-Phage Coordination"
    add_header(slide, prs, title)

    # 左侧内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    if lang == 'cn':
        bullets = [
            "噬菌体多样性显著降低 (P=0.032)",
            "",
            "细菌-噬菌体高度协同：",
            "  • Mantel r=0.52, P=0.001",
            "  • Procrustes P=0.002",
            "",
            "三角网络：17,848条边",
            "",
            "级联效应：",
            "  CagA → 细菌 → 噬菌体"
        ]
    else:
        bullets = [
            "Phage diversity significantly decreased (P=0.032)",
            "",
            "Bacteria-phage highly coordinated:",
            "  • Mantel r=0.52, P=0.001",
            "  • Procrustes P=0.002",
            "",
            "Tripartite network: 17,848 edges",
            "",
            "Cascade effect:",
            "  CagA → Bacteria → Phage"
        ]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        if "P=0.001" in bullet or "P=0.002" in bullet or "P=0.032" in bullet:
            p.font.bold = True
            p.font.color.rgb = COLORS['green']
        elif "级联" in bullet or "Cascade" in bullet:
            p.font.bold = True
            p.font.color.rgb = COLORS['primary_blue']
        else:
            p.font.color.rgb = COLORS['dark_gray']

    # 右侧：Procrustes图或级联效应示意图
    cascade_title = slide.shapes.add_textbox(Inches(5.0), Inches(1.3), Inches(4.5), Inches(0.4))
    tf = cascade_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Mantel & Procrustes分析" if lang == 'cn' else "Mantel & Procrustes Analysis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    if image_path and os.path.exists(image_path):
        # 使用Procrustes图片
        slide.shapes.add_picture(image_path, Inches(5.0), Inches(1.7), width=Inches(4.7))
    else:
        # 备用：级联效应示意图
        add_box_with_text(slide, 6.2, 1.9, 1.5, 0.6, "CagA\n感染" if lang == 'cn' else "CagA\nInfection", COLORS['red'], COLORS['white'], 12, True)

        arr1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.85), Inches(2.55), Inches(0.2), Inches(0.3))
        arr1.fill.solid()
        arr1.fill.fore_color.rgb = COLORS['dark_gray']
        arr1.line.fill.background()

        add_box_with_text(slide, 6.2, 2.9, 1.5, 0.6, "细菌\n重塑" if lang == 'cn' else "Bacteria\nRestructured", COLORS['primary_blue'], COLORS['white'], 12, True)

        arr2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.85), Inches(3.55), Inches(0.2), Inches(0.3))
        arr2.fill.solid()
        arr2.fill.fore_color.rgb = COLORS['dark_gray']
        arr2.line.fill.background()

        add_box_with_text(slide, 6.2, 3.9, 1.5, 0.6, "噬菌体\n协同变化" if lang == 'cn' else "Phage\nCoordinated", COLORS['purple'], COLORS['white'], 12, True)

    if lang == 'cn':
        notes = """最后一个结果是细菌和噬菌体的协同变化。

前面我们说CagA不直接影响病毒组，但这不意味着病毒组没有变化。恰恰相反，我们发现噬菌体也发生了显著改变——只是这种改变是细菌变化的下游效应。

左侧是关键统计结果。首先，噬菌体的Alpha多样性显著降低，P值为0.032。这说明CagA感染后，噬菌体的种类减少了。

更重要的是细菌-噬菌体的协同分析。我们做了两种分析：Mantel检验和Procrustes分析。Mantel检验r值达到0.52，P值仅为0.001，高度显著；Procrustes分析P值为0.002，同样高度显著。这说明细菌群落和噬菌体群落的结构变化是高度一致的。

右图展示了这两种分析的结果。左边是Alpha多样性相关图，虽然相关性不显著，说明多样性水平不是关联的关键。右边是Procrustes分析图，你可以看到细菌和噬菌体的样本点之间用线连接，线越短说明两者越一致。不同颜色代表不同处理组。

结论是什么？CagA的效应通过级联传递：CagA首先改变细菌群落，细菌的变化又带动了噬菌体的改变。这是一个完整的微生物组级联效应。"""
    else:
        notes = """The final result is bacteria-phage coordination.

We said CagA doesn't directly affect viruses, but this doesn't mean viruses are unchanged. On the contrary, we found significant phage changes — they're just downstream effects of bacterial changes.

The left shows key statistics. First, phage Alpha diversity significantly decreased with P=0.032. This means after CagA infection, phage species diversity decreased.

More importantly, bacteria-phage coordination analysis. We performed two analyses: Mantel test and Procrustes analysis. Mantel test r=0.52 with P=0.001, highly significant; Procrustes P=0.002, also highly significant. This means bacterial and phage community structural changes are highly consistent.

The right figure shows these two analyses. Left is Alpha diversity correlation — not significant, meaning diversity levels aren't the key association. Right is Procrustes analysis — you can see bacterial and phage sample points connected by lines, shorter lines mean better agreement. Different colors represent different treatment groups.

What's the conclusion? CagA's effect propagates through a cascade: CagA first changes bacteria, then bacterial changes drive phage changes. This is a complete microbiome cascade effect."""

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

    return slide


def add_conclusion_slide(prs, title, conclusions, notes=""):
    """添加结论幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, prs, title)

    for i, (num, text) in enumerate(conclusions):
        y_pos = 1.35 + i * 0.75

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['green']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.03), Inches(0.45), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_pos), Inches(8.4), Inches(0.65))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['dark_gray']

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_thanks_slide(prs, text, subtext=""):
    """添加致谢幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_blue']
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_chinese_ppt(output_path, image_dir):
    """生成中文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. 封面
    add_title_slide(
        prs,
        "CagA依赖性肠道微生物组重塑\n及其功能冗余特征",
        "PC047 vCagAepitope 微生物组分析进展",
        "2026-01-26",
        "汇报人：龚宇航"
    )

    # 2. 背景1：CagA与肠道肿瘤
    add_background_slide_1(prs, 'cn')

    # 3. 背景2：Apc突变与G×E
    add_background_slide_2(prs, 'cn')

    # 4. 核心假说图解
    add_hypothesis_slide(prs, 'cn')

    # 5. 分析流程图
    add_analysis_pipeline_slide(prs, 'cn')

    # 6. 结果1：CagA显著重塑肠道菌群（使用综合Summary图）
    summary_figure_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "47_part4_summary_figure.png")
    add_content_slide(
        prs,
        "结果1：CagA显著重塑肠道菌群",
        [
            "• Beta多样性：**P=0.012**, R²=33.6%",
            "• G×E交互：**P=0.024** (显著)",
            "• CagA效应@ApcMUT：**显著**",
            "• CagA效应@ApcWT：不显著",
        ],
        notes="""现在来看第一个核心结果：CagA显著重塑肠道菌群。

先看左侧的要点总结。Beta多样性分析显示P值仅为0.012，效应量R²达到33.6%，说明CagA解释了三分之一的菌群变异。同时，G×E交互作用也是显著的，P值0.024。最重要的是，CagA效应在Apc突变背景下显著，但在野生型背景下不显著。

右图是我们的综合分析结果，信息量很大。

左上角是PCoA图，展示了全部6组样本的群落分布。你可以看到不同颜色代表不同的感染状态——红色是对照组，蓝色是HpKO感染组，绿色是HpWT也就是CagA阳性感染组。椭圆表示95%置信区间。可以看到各组之间有明显的分离。

右上角是效应量图，展示了2×3因子设计中各因素解释的方差比例。最关键的是交互作用项，占到了100%，说明CagA的效应确实依赖于遗传背景。

下方的条形图是成对比较的结果。你可以看到蓝色条代表CagA效应，在Apc突变小鼠中P值小于0.05，是显著的；但在野生型小鼠中就不显著了。

这是典型的G×E交互作用模式，说明CagA的菌群重塑效应需要Apc突变这个"易感背景"。""",
        image_path=summary_figure_path if os.path.exists(summary_figure_path) else None
    )

    # 7. 结果2：细菌vs病毒
    add_bacteria_vs_virus_slide(prs, 'cn')

    # 8. 结果3：功能冗余全层级验证（新增GO/PFAM）+ Driver图
    driver_shift_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_functional_redundancy_slide(prs, 'cn', image_path=driver_shift_path if os.path.exists(driver_shift_path) else None)

    # 9. 结果4：共现网络分析（新增）
    add_network_analysis_slide(prs, 'cn')

    # 10. 结果5：噬菌体协同变化 + Procrustes图
    procrustes_path = os.path.join(image_dir, "03_singlem_diversity_analysis", "28_host_virus_association.png")
    add_phage_coordination_slide(prs, 'cn', image_path=procrustes_path if os.path.exists(procrustes_path) else None)

    # 11. 结论
    add_conclusion_slide(
        prs,
        "核心结论与后续计划",
        [
            (1, "CagA显著重塑肠道菌群：Beta多样性P=0.012，选择性调控"),
            (2, "功能冗余全层级验证：Pathway→KO→GO→PFAM，FDR均无显著"),
            (3, "网络模块结构瓦解：模块度降低62%，解释功能冗余机制"),
            (4, "Driver静默更替：益生菌↓致病菌↑，改变免疫刺激谱"),
            (5, "细菌-噬菌体协同：Mantel P=0.001，级联效应"),
            (6, "后续：代谢组学验证 + 整合肿瘤/T细胞表型数据"),
        ],
        notes="""好的，现在让我来总结一下今天汇报的核心结论。

第一，CagA显著重塑肠道菌群结构。Beta多样性P值仅0.012，效应量33.6%。重要的是，这是选择性调控而非全面破坏——Alpha多样性没有变化。

第二，功能冗余得到全层级验证。从Pathway到KO到GO再到PFAM，四个功能层级FDR校正后均无显著差异。虽然物种变了，但功能稳定。

第三，网络模块结构瓦解。共现网络分析显示模块度降低62%，从0.468降到0.177。这从网络层面解释了功能冗余的机制——物种间的功能分工被打乱了。

第四，Driver物种发生了静默更替。虽然功能总量不变，但执行功能的物种已经换了——从益生菌转向致病菌。L. johnsonii下降50%，M. schaedleri上升57%。这对免疫系统来说是完全不同的信号。

第五，细菌和噬菌体高度协同变化。Mantel检验P值仅0.001，存在明显的级联效应。

后续计划有两个方向：一是代谢组学验证，看看功能冗余表象下代谢物是否真的相同；二是整合肿瘤和T细胞表型数据，验证假说的后半部分。

感谢大家的聆听！欢迎提问和讨论。"""
    )

    # 12. 致谢
    add_thanks_slide(
        prs,
        "感谢聆听",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"中文版PPT已生成：{output_path}")


def generate_english_ppt(output_path, image_dir):
    """生成英文版PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title
    add_title_slide(
        prs,
        "CagA-Dependent Gut Microbiome Restructuring\nand Functional Redundancy",
        "PC047 vCagAepitope Microbiome Analysis Progress",
        "2026-01-26",
        "Presenter: Yuhang Gong"
    )

    # 2. Background 1
    add_background_slide_1(prs, 'en')

    # 3. Background 2
    add_background_slide_2(prs, 'en')

    # 4. Hypothesis
    add_hypothesis_slide(prs, 'en')

    # 5. Pipeline
    add_analysis_pipeline_slide(prs, 'en')

    # 6. Result 1: CagA Restructures Microbiota（使用综合Summary图）
    summary_figure_path = os.path.join(image_dir, "01_alpha_beta_diversity_analysis", "47_part4_summary_figure.png")
    add_content_slide(
        prs,
        "Result 1: CagA Restructures Gut Microbiota",
        [
            "• Beta diversity: **P=0.012**, R²=33.6%",
            "• G×E interaction: **P=0.024** (significant)",
            "• CagA effect@ApcMUT: **significant**",
            "• CagA effect@ApcWT: not significant",
        ],
        notes="""Now let's look at the first core result: CagA significantly restructures gut microbiota.

First, the key points on the left. Beta diversity analysis shows P=0.012, with effect size R² reaching 33.6% — meaning CagA explains one-third of microbiome variation. G×E interaction is also significant at P=0.024. Most importantly, CagA effect is significant under Apc-mutant background but not significant under wild-type background.

The figure on the right shows our comprehensive analysis — very information-rich.

Top-left is the PCoA plot showing community distribution of all 6 groups. Different colors represent different infection states — red for control, blue for HpKO infection, green for HpWT or CagA-positive infection. Ellipses show 95% confidence intervals. You can see clear separation between groups.

Top-right is the effect size plot showing variance explained by each factor in the 2×3 factorial design. Most critically, the interaction term accounts for 100%, indicating CagA's effect indeed depends on genetic background.

The bottom bar chart shows pairwise comparison results. Blue bars represent CagA effect — significant in Apc-mutant mice with P<0.05, but not significant in wild-type mice.

This is a classic G×E interaction pattern, showing CagA's microbiome restructuring effect requires the "susceptible background" of Apc mutation.""",
        image_path=summary_figure_path if os.path.exists(summary_figure_path) else None
    )

    # 7. Result 2: Bacteria vs Virus
    add_bacteria_vs_virus_slide(prs, 'en')

    # 8. Result 3: Functional redundancy (with GO/PFAM) + Driver image
    driver_shift_path = os.path.join(image_dir, "02_functional_profiling", "91_part9_driver_species_shift.png")
    add_functional_redundancy_slide(prs, 'en', image_path=driver_shift_path if os.path.exists(driver_shift_path) else None)

    # 9. Result 4: Network analysis (new)
    add_network_analysis_slide(prs, 'en')

    # 10. Result 5: Phage coordination + Procrustes image
    procrustes_path = os.path.join(image_dir, "03_singlem_diversity_analysis", "28_host_virus_association.png")
    add_phage_coordination_slide(prs, 'en', image_path=procrustes_path if os.path.exists(procrustes_path) else None)

    # 11. Conclusions
    add_conclusion_slide(
        prs,
        "Conclusions & Future Directions",
        [
            (1, "CagA restructures gut microbiota: Beta P=0.012, selective regulation"),
            (2, "Functional redundancy at all levels: Pathway→KO→GO→PFAM, no FDR significance"),
            (3, "Network module collapse: Modularity -62%, explains redundancy mechanism"),
            (4, "Silent Driver shift: Probiotics↓ Pathobionts↑, altering immune stimulation"),
            (5, "Bacteria-phage coordination: Mantel P=0.001, cascade effect"),
            (6, "Next: Metabolomics validation + Tumor/T-cell phenotype integration"),
        ],
        notes="""Alright, let me summarize today's core conclusions.

First, CagA significantly restructures gut microbiota. Beta diversity P=0.012, effect size 33.6%. Importantly, this is selective regulation, not destruction — Alpha diversity remained unchanged.

Second, functional redundancy verified at all levels. From Pathway to KO to GO to PFAM, all four functional levels showed no FDR-significant differences. Species changed, but functions remained stable.

Third, network module structure collapsed. Co-occurrence network analysis showed modularity decreased 62%, from 0.468 to 0.177. This explains functional redundancy at the network level — functional division among species was disrupted.

Fourth, Driver species silently shifted. Although total function remained unchanged, the species executing functions changed — from probiotics to pathobionts. L. johnsonii decreased 50%, M. schaedleri increased 57%. For the immune system, this is a completely different signal.

Fifth, bacteria and phage are highly coordinated. Mantel test P=0.001, showing clear cascade effect.

Future directions: First, metabolomics validation to see if metabolites are truly the same under functional redundancy; second, integration with tumor and T-cell phenotype data to verify the latter half of our hypothesis.

Thank you for your attention! Questions and discussions are welcome."""
    )

    # 12. Thanks
    add_thanks_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    add_page_numbers_to_presentation(prs, skip_first=False, skip_last=False)

    prs.save(output_path)
    print(f"英文版PPT已生成：{output_path}")


if __name__ == "__main__":
    base_dir = r"C:\Users\36094\Desktop\pc047oma"
    image_dir = os.path.join(base_dir, "analyses", "data")

    # 汇报日期和版本号
    meeting_date = "20260126"
    version = "v6"

    # 创建日期文件夹
    output_dir = os.path.join(base_dir, "PPT", meeting_date)
    os.makedirs(output_dir, exist_ok=True)

    # 输出路径
    chinese_path = os.path.join(output_dir, f"组会汇报_{version}.pptx")
    english_path = os.path.join(output_dir, f"GroupMeeting_{version}.pptx")

    generate_chinese_ppt(chinese_path, image_dir)
    generate_english_ppt(english_path, image_dir)

    print("\n生成完成！")
    print(f"中文版：{chinese_path}")
    print(f"英文版：{english_path}")
